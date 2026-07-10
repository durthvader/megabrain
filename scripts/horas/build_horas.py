"""Monta o PPTX executivo de Custo de Horas a partir do JSON + narrativa.

Camada 2 da máquina (montador). Consome:
  - outputs/horas/dados_<regiao>.json   (números — gerado pelo extrair_horas.py)
  - docs/horas/narrativa_<regiao>.json  (textos — títulos, causas, ações)
e gera outputs/horas/Horas_<Regiao>_Executivo.pptx no padrão Alloha.

Dinâmica dos 5 slides: Fato (indicador) -> Causa/onde (ociosidade + regimes
útil-noturno x FDS-diurno) -> Causa/quem (GA estrutural x eventual) ->
Risco x custo-benefício (dimensionamento seguro + conta da conversão) ->
Ação (proposta do próximo ciclo + FCA cirúrgico).

Uso:
    python scripts/horas/build_horas.py --regiao ceara
"""

from __future__ import annotations

import argparse
import json
import unicodedata
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[2]

# ------------------------------------------------------------ identidade Alloha
NAVY = "081551"
BLUE = "1476C9"
BLUE_MID = "3E6FB0"
BLUE_LIGHT = "DCECF9"
GREEN = "42DF4B"
GREEN_DARK = "13853A"
RED = "D94545"
RED_LIGHT = "FDE2DF"
ORANGE = "F08A24"
AMBER = "E4BD30"
WARN = "9A6200"
INK = "182238"
MUTED = "667085"
LINE = "D9E0EA"
WASH = "F3F6FA"
WHITE = "FFFFFF"

MES_ABREV = {"janeiro": "Jan", "fevereiro": "Fev", "março": "Mar", "abril": "Abr",
             "maio": "Mai", "junho": "Jun", "julho": "Jul", "agosto": "Ago",
             "setembro": "Set", "outubro": "Out", "novembro": "Nov", "dezembro": "Dez"}


# ------------------------------------------------------------------- primitivas
def px(value: float) -> Emu:
    return Emu(int(value * 9525))


def rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


def add_rect(slide, left, top, width, height, fill, line=None, line_width=1.0, round_corners=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_corners else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, px(left), px(top), px(width), px(height))
    if round_corners:
        shape.adjustments[0] = 0.06
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_text(slide, value, left, top, width, height, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Aptos", line_spacing=1.06):
    box = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, margin, 0)
    lines = value if isinstance(value, list) else [(value, {})]
    for index, (text, opts) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = opts.get("align", align)
        paragraph.line_spacing = opts.get("line_spacing", line_spacing)
        if index > 0:
            paragraph.space_before = Pt(opts.get("space_before", 2))
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        run.font.color.rgb = rgb(opts.get("color", color))
        run.font.name = opts.get("font", font)
    return box


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def chrome(slide, number, total, source):
    add_rect(slide, 0, 0, 1100, 8, NAVY)
    add_rect(slide, 1100, 0, 180, 8, GREEN)
    add_text(slide, "alloha", 1090, 30, 130, 36, size=27, color=NAVY,
             align=PP_ALIGN.RIGHT, font="Aptos Display")
    add_text(slide, "F I B R A", 1090, 66, 130, 14, size=8, color=GREEN_DARK,
             bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, source, 56, 690, 1040, 16, size=9, color="8B94A6")
    add_text(slide, f"{number} / {total}", 1150, 690, 74, 16, size=9,
             color="8B94A6", align=PP_ALIGN.RIGHT)


def title_block(slide, eyebrow, title, lead=None, title_size=30, lead_top=None):
    add_text(slide, eyebrow, 56, 34, 980, 20, size=12, color=BLUE, bold=True)
    add_text(slide, title, 56, 58, 1050, 76, size=title_size, color=NAVY,
             bold=True, font="Aptos Display", line_spacing=1.0)
    if lead:
        add_text(slide, lead, 56, lead_top or 128, 1160, 46, size=14.5, color="45516A")


def set_cell(cell, lines, fill=None, anchor=MSO_ANCHOR.MIDDLE, v_margin=5):
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill)
    cell.vertical_anchor = anchor
    cell.margin_left = px(9)
    cell.margin_right = px(8)
    cell.margin_top = px(v_margin)
    cell.margin_bottom = px(v_margin)
    tf = cell.text_frame
    tf.word_wrap = True
    for index, (text, opts) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = opts.get("align", PP_ALIGN.LEFT)
        paragraph.line_spacing = opts.get("line_spacing", 1.04)
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", 11))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = rgb(opts.get("color", INK))
        run.font.name = opts.get("font", "Aptos")


def add_table(slide, rows, cols, left, top, width, height, col_widths):
    frame = slide.shapes.add_table(rows, cols, px(left), px(top), px(width), px(height))
    table = frame.table
    table.first_row = False
    table.horz_banding = False
    for index, cw in enumerate(col_widths):
        table.columns[index].width = px(cw)
    return table


def hbar(slide, label, value_text, left, top, label_width, bar_width, ratio,
         color=BLUE, track="E7EDF5", bar_height=13, val_color=INK):
    add_text(slide, label, left, top - 2, label_width, 16, size=10.5, color=INK)
    track_left = left + label_width + 8
    add_rect(slide, track_left, top, bar_width, bar_height, track)
    add_rect(slide, track_left, top, max(bar_width * ratio, 3), bar_height, color)
    add_text(slide, value_text, track_left + bar_width + 8, top - 2, 70, 16,
             size=11, color=val_color, bold=True)


# --------------------------------------------------------------------- formato
def moeda(v):
    if v is None:
        return "—"
    return "R$ " + f"{v / 1000:.1f}".replace(".", ",") + " mil"


def pct(x, casas=0):
    if x is None:
        return "—"
    return f"{x * 100:.{casas}f}".replace(".", ",") + "%"


def encurtar_nome(nome, n=2):
    if not nome:
        return "—"
    partes = str(nome).strip().split()
    sel = partes[:n]
    conect = {"da", "de", "do", "das", "dos", "e"}
    while len(sel) > 1 and sel[-1].lower() in conect:
        sel = sel[:-1]
    return " ".join(sel).title()


def cesta_get(cesta, *chaves):
    """Casa uma chave da cesta ignorando acento/caixa."""
    def nrm(s):
        s = str(s).lower()
        return "".join(c for c in unicodedata.normalize("NFKD", s) if not unicodedata.combining(c))
    alvo = [nrm(c) for c in chaves]
    for k, v in cesta.items():
        if nrm(k) in alvo:
            return v
    return None


class SafeDict(dict):
    def __missing__(self, key):
        return "{" + key + "}"


def fmt(texto, ctx):
    return str(texto).format_map(SafeDict(ctx))


# ---------------------------------------------------------------------- charts
def line_chart(slide, series, categorias, left, top, width, height):
    """Gráfico de linhas nativo (valores em R$ mil). series = [(nome, [vals], cor, largura)]."""
    cd = CategoryChartData()
    cd.categories = categorias
    for nome, vals, _c, _w in series:
        cd.add_series(nome, [None if v is None else v / 1000.0 for v in vals])
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, px(left), px(top),
                                     px(width), px(height), cd)
    chart = graphic.chart
    chart.has_title = False
    chart.font.name = "Aptos"
    chart.font.size = Pt(9)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    for s, (nome, vals, cor, largura) in zip(chart.series, series):
        s.format.line.color.rgb = rgb(cor)
        s.format.line.width = Pt(largura)
        s.smooth = False
        s.marker.style = 8 if nome.lower().startswith("real") else 2  # circle p/ real
        s.marker.size = 5
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = rgb(cor)
        s.marker.format.line.color.rgb = rgb(cor)
    cat = chart.category_axis
    cat.tick_labels.font.size = Pt(9)
    cat.has_major_gridlines = False
    cat.major_tick_mark = XL_TICK_MARK.OUTSIDE
    cat.format.line.color.rgb = rgb(LINE)
    val = chart.value_axis
    val.has_major_gridlines = True
    val.major_gridlines.format.line.color.rgb = rgb("EDF1F6")
    val.tick_labels.font.size = Pt(9)
    val.tick_labels.number_format = '#,##0'
    val.tick_labels.number_format_is_linked = False
    val.format.line.fill.background()
    return chart


def stacked_horizontal(slide, itens, left, top, width, height):
    """Barra horizontal empilhada. itens = [(valor, cor)]. Retorna limites de cada segmento."""
    total = sum(v for v, _ in itens) or 1
    x = left
    limites = []
    for valor, cor in itens:
        w = width * valor / total
        add_rect(slide, x, top, w, height, cor)
        limites.append((x, w))
        x += w
    return limites


# ---------------------------------------------------------------------- slides
def slide1(prs, dados, nar, ctx, blank):
    slide = prs.slides.add_slide(blank)
    s = nar["s1"]
    chrome(slide, 1, 5, fmt(nar["fonte"], ctx))
    title_block(slide, fmt(s["eyebrow"], ctx), fmt(s["titulo"], ctx), title_size=28)

    # --- três KPIs
    cards = [
        ("ORÇADO DO CICLO", ctx["orcado"], "por mês", False),
        ("META · −50%", ctx["meta"], "compromisso do próximo ciclo", True),
        ("REALIZADO (FOTO)", ctx["real"], f"{ctx['var_orcado']} vs orçado · falta cortar {ctx['gap_pct']}", False),
    ]
    cx = 56
    for titulo, valor, sub, destaque in cards:
        w = 372
        fill = NAVY if destaque else WHITE
        add_rect(slide, cx, 178, w, 96, fill, line=None if destaque else LINE, round_corners=True)
        add_text(slide, titulo, cx + 18, 192, w - 36, 16, size=10.5,
                 color="9FC9EE" if destaque else BLUE, bold=True)
        add_text(slide, valor, cx + 18, 210, w - 36, 40, size=30, bold=True,
                 color=WHITE if destaque else NAVY, font="Aptos Display")
        add_text(slide, sub, cx + 18, 250, w - 36, 18, size=10,
                 color="C9D6EF" if destaque else MUTED)
        cx += w + 20

    # --- gráfico de tendência (esquerda)
    serie = [m for m in dados["serie_total"] if m["real"] is not None]
    cats = [MES_ABREV.get(m["mes"], m["mes"][:3]) for m in serie]
    reais = [m["real"] for m in serie]
    orcs = [m["orcado"] for m in serie]
    metas = [ctx["meta_val"]] * len(serie)
    add_text(slide, "Custo mensal de horas — R$ mil", 56, 296, 560, 18, size=12.5,
             color=NAVY, bold=True)
    line_chart(slide, [("Realizado", reais, NAVY, 2.5),
                       ("Orçado", orcs, MUTED, 1.75),
                       ("Meta −50%", metas, GREEN_DARK, 1.75)],
               cats, 44, 316, 585, 336)

    # --- composição da cesta (direita)
    add_text(slide, "Composição do custo (mês-foto)", 660, 296, 560, 18, size=12.5,
             color=NAVY, bold=True)
    ordem = [("DSR", cesta_get(dados["cesta"], "DSR"), BLUE),
             ("Acionamento sobreaviso", cesta_get(dados["cesta"], "ACIONAMENTO SOBRE AVISO"), NAVY),
             ("Espera sobreaviso", cesta_get(dados["cesta"], "ESPERA SOBRE AVISO"), BLUE_MID),
             ("Banco de horas", cesta_get(dados["cesta"], "BH"), "98A2B5"),
             ("Feriados", cesta_get(dados["cesta"], "FERIADOS"), "C9D6E6")]
    ordem = [(n, v or 0, c) for n, v, c in ordem]
    total = sum(v for _, v, _ in ordem) or 1
    limites = stacked_horizontal(slide, [(v, c) for _, v, c in ordem], 660, 320, 562, 34)

    # colchete "Sobreaviso" SOB Acionamento + Espera — abaixo da barra p/ nunca
    # colidir com o título, independente do tamanho do DSR
    x_ac = limites[1][0]
    w_se = limites[1][1] + limites[2][1]
    add_rect(slide, x_ac, 358, w_se, 3, GREEN_DARK)
    add_text(slide, f"Sobreaviso = {ctx['sobreaviso_pct']}",
             x_ac, 362, w_se, 16, size=10.5, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

    # legenda
    ly = 388
    for nome, valor, cor in ordem:
        add_rect(slide, 660, ly + 2, 12, 12, cor)
        add_text(slide, nome, 680, ly, 300, 16, size=10.5, color=INK)
        add_text(slide, f"{moeda(valor)}  ·  {pct(valor / total)}", 680, ly, 552, 16,
                 size=10.5, color=MUTED, align=PP_ALIGN.RIGHT)
        ly += 22

    # decisão
    add_rect(slide, 660, 512, 562, 128, WASH, line=LINE, round_corners=True)
    add_text(slide, "DECISÃO", 680, 524, 300, 16, size=10.5, color=BLUE, bold=True)
    add_text(slide, fmt(s["decisao"], ctx), 680, 544, 524, 90, size=13, color=INK, bold=True,
             line_spacing=1.08)

    add_notes(slide, "SLIDE 1 — INDICADOR (Fato). Números vêm de dados_<regiao>.json; "
                     "textos de narrativa_<regiao>.json (s1). Meta = -50% do orçado.")


def slide2(prs, dados, nar, ctx, blank):
    slide = prs.slides.add_slide(blank)
    s = nar["s2"]
    o = dados["ociosidade"]
    chrome(slide, 2, 5, fmt(nar["fonte"], ctx))
    title_block(slide, fmt(s["eyebrow"], ctx), fmt(s["titulo"], ctx), title_size=27)

    # painel A — ociosidade por headcount
    add_rect(slide, 56, 150, 560, 250, WHITE, line=LINE, round_corners=True)
    add_text(slide, "Técnicos por dia: em espera x acionados", 78, 166, 520, 18,
             size=13, color=NAVY, bold=True)
    add_text(slide, "Média por dia no período de sobreaviso", 78, 188, 520, 14,
             size=10, color=MUTED)
    he, ha = o["hc_espera_medio"], o["hc_acionado_medio"]
    hbar(slide, "Em espera", f"{he:.1f}".replace(".", ",").rstrip("0").rstrip(","),
         78, 232, 130, 300, 1.0, color=BLUE_MID)
    hbar(slide, "Acionados", f"{ha:.1f}".replace(".", ",").rstrip("0").rstrip(","),
         78, 268, 130, 300, ha / he if he else 0, color=NAVY)
    add_rect(slide, 78, 312, 520, 66, WASH, round_corners=True)
    add_text(slide, ctx["ocio_pct_hc"], 96, 322, 150, 46, size=30, color=GREEN_DARK,
             bold=True, font="Aptos Display", anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, "dos técnicos de sobreaviso são de fato acionados por dia — "
                    "o restante é plantão pago em espera.",
             250, 322, 332, 48, size=11.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # painel B — ociosidade em horas
    add_rect(slide, 636, 150, 586, 250, WHITE, line=LINE, round_corners=True)
    add_text(slide, "Horas de sobreaviso: espera x acionamento", 658, 166, 540, 18,
             size=13, color=NAVY, bold=True)
    add_text(slide, "Soma do período (analítico nominal)", 658, 188, 540, 14,
             size=10, color=MUTED)
    esp, aci = o["horas_espera_total"], o["horas_acionamento_total"]
    tot = esp + aci
    add_text(slide, "Espera", 658, 210, 120, 16, size=10.5, color=INK)
    add_rect(slide, 658, 228, 540, 26, "E7EDF5")
    add_rect(slide, 658, 228, 540, 26, BLUE_MID)
    add_text(slide, f"{esp:,.0f} h".replace(",", "."), 658, 231, 528, 18, size=11,
             color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "Acionamento efetivo", 658, 266, 200, 16, size=10.5, color=INK)
    add_rect(slide, 658, 284, 540, 26, "E7EDF5")
    add_rect(slide, 658, 284, max(540 * aci / tot, 6), 26, NAVY)
    add_text(slide, f"{aci:,.0f} h".replace(",", "."), 658 + max(540 * aci / tot, 6) + 8,
             287, 120, 18, size=11, color=NAVY, bold=True)
    add_rect(slide, 658, 332, 540, 54, NAVY, round_corners=True)
    add_text(slide, ctx["ocio_pct_horas"], 676, 332, 96, 54, size=24, color=GREEN,
             bold=True, font="Aptos Display", anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, "do custo de sobreaviso vira acionamento — o resto é plantão em espera.",
             776, 332, 406, 54, size=11.5, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # painel C — dois regimes: dia útil x fim de semana, por faixa horária
    intel = dados.get("inteligencia") or {}
    reg = intel.get("regimes") or {}
    add_rect(slide, 56, 416, 1166, 218, WHITE, line=LINE, round_corners=True)
    add_text(slide, "Não é um problema só: no dia útil o acionamento é noturno; no fim de semana, diurno",
             78, 430, 900, 18, size=13, color=NAVY, bold=True)
    add_text(slide, "Horas de acionamento por faixa horária, separadas em dia útil × sábado/domingo.",
             78, 452, 900, 16, size=10.5, color=MUTED)
    faixas_nomes = ["Noite", "Madrugada", "Dia", "Tarde"]
    maxf = max((reg.get(f, {}).get("util", 0) + reg.get(f, {}).get("fds", 0))
               for f in faixas_nomes) or 1
    fy = 486
    for f in faixas_nomes:
        u = reg.get(f, {}).get("util", 0)
        w = reg.get(f, {}).get("fds", 0)
        tot_f = u + w
        add_text(slide, f, 78, fy - 2, 92, 16, size=10.5, color=INK)
        bar_l, bar_w = 174, 440
        add_rect(slide, bar_l, fy, bar_w, 14, "E7EDF5")
        wu = bar_w * u / maxf
        ww = bar_w * w / maxf
        if wu > 0:
            add_rect(slide, bar_l, fy, max(wu, 2), 14, NAVY)
        if ww > 0:
            add_rect(slide, bar_l + wu, fy, max(ww, 2), 14, GREEN_DARK)
        pf = (w / tot_f) if tot_f else 0
        add_text(slide, f"{tot_f:.0f} h", bar_l + bar_w + 10, fy - 2, 64, 16,
                 size=10.5, color=INK, bold=True)
        add_text(slide, f"{pct(pf)} FDS", bar_l + bar_w + 76, fy - 2, 80, 16, size=10.5,
                 color=GREEN_DARK if pf >= 0.5 else MUTED, bold=pf >= 0.5)
        fy += 34
    # legenda mini
    add_rect(slide, 174, fy + 2, 11, 11, NAVY)
    add_text(slide, "dia útil", 190, fy, 70, 14, size=9.5, color=MUTED)
    add_rect(slide, 254, fy + 2, 11, 11, GREEN_DARK)
    add_text(slide, "fim de semana", 270, fy, 110, 14, size=9.5, color=MUTED)

    # insight lateral — dois regimes
    add_rect(slide, 850, 480, 348, 150, WASH, round_corners=True)
    add_text(slide, "DOIS REGIMES, DUAS SOLUÇÕES", 870, 490, 310, 14, size=10, color=BLUE, bold=True)
    add_text(slide, fmt(s["insight"], ctx), 870, 508, 312, 116, size=10, color=INK,
             line_spacing=1.07)

    add_notes(slide, "SLIDE 2 — CAUSA/onde (ociosidade + regimes). HC/dia e horas de espera x "
                     "acionamento; painel C mostra dia útil (noturno) x FDS (diurno). "
                     "Números do JSON (ociosidade, inteligencia.regimes).")


def slide3(prs, dados, nar, ctx, blank):
    slide = prs.slides.add_slide(blank)
    s = nar["s3"]
    chrome(slide, 3, 5, fmt(nar["fonte"], ctx))
    title_block(slide, fmt(s["eyebrow"], ctx), fmt(s["titulo"], ctx), title_size=27)

    gas = dados["sobreaviso_por_ga"][:8]
    maxt = max(g["horas_total"] for g in gas) or 1

    # ranking à esquerda (barras empilhadas espera+acionamento)
    add_text(slide, "Horas de sobreaviso por Área de Gestão (GA)", 56, 152, 620, 18,
             size=13, color=NAVY, bold=True)
    add_text(slide, "Barra = espera + acionamento · maior no topo",
             56, 174, 620, 14, size=10, color=MUTED)
    y = 202
    barmax = 380
    for g in gas:
        add_text(slide, encurtar_nome(g["gestor"]), 56, y - 2, 170, 16, size=10.5, color=INK)
        w_tot = barmax * g["horas_total"] / maxt
        w_esp = w_tot * (g["horas_espera"] / g["horas_total"] if g["horas_total"] else 0)
        add_rect(slide, 232, y, max(w_tot, 3), 15, BLUE_MID)
        add_rect(slide, 232, y, max(w_esp, 2), 15, NAVY)
        add_text(slide, f"{g['horas_total']:.0f} h", 232 + max(w_tot, 3) + 8, y - 2, 70, 16,
                 size=10.5, color=INK, bold=True)
        y += 30

    # legenda mini
    add_rect(slide, 232, y + 4, 11, 11, NAVY)
    add_text(slide, "espera", 248, y + 2, 70, 14, size=9.5, color=MUTED)
    add_rect(slide, 312, y + 4, 11, 11, BLUE_MID)
    add_text(slide, "acionamento", 328, y + 2, 100, 14, size=9.5, color=MUTED)

    # painel direito — Pareto + tabela
    add_rect(slide, 700, 150, 522, 150, NAVY, round_corners=True)
    add_text(slide, "CONCENTRAÇÃO", 722, 166, 300, 14, size=10.5, color="9FC9EE", bold=True)
    add_text(slide, ctx["top3_ga_pct"], 722, 186, 160, 48, size=38, color=GREEN,
             bold=True, font="Aptos Display")
    add_text(slide, fmt(s["callout"], ctx), 722, 238, 480, 56, size=11, color=WHITE,
             line_spacing=1.06)

    # tabela por GA: frequência de acionamento define o instrumento certo
    intel = dados.get("inteligencia") or {}
    gas_i = (intel.get("gas") or [])[:6]
    classes = {"estrutural": ("TODO DIA", RED),
               "intermediario": ("FREQUENTE", WARN),
               "eventual": ("EVENTUAL", GREEN_DARK)}
    table = add_table(slide, len(gas_i) + 1, 4, 700, 320, 522, 300, [204, 112, 96, 110])
    cabec = ["Área de gestão (GA)", "Aciona em", "Custo", "Padrão"]
    for c, h in enumerate(cabec):
        set_cell(table.cell(0, c), [(h, {"size": 10.5, "bold": True, "color": WHITE})],
                 fill=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    for r, g in enumerate(gas_i, start=1):
        rotulo, cor = classes.get(g["classe"], ("—", MUTED))
        fill = WHITE if r % 2 else WASH
        set_cell(table.cell(r, 0), [(encurtar_nome(g["gestor"], 3), {"size": 10.5})], fill=fill)
        set_cell(table.cell(r, 1), [(f"{g['dias_acionamento']} de {g['dias_escala']} d",
                 {"size": 10.5, "align": PP_ALIGN.CENTER, "bold": True})], fill=fill)
        set_cell(table.cell(r, 2), [(moeda(g["custo_total"]), {"size": 10.5, "align": PP_ALIGN.CENTER})], fill=fill)
        set_cell(table.cell(r, 3), [(rotulo, {"size": 9.5, "align": PP_ALIGN.CENTER,
                 "bold": True, "color": cor})], fill=fill)

    add_notes(slide, "SLIDE 3 — CAUSA/quem (por GA). Ranking de horas + frequência de acionamento: "
                     "GA que aciona todo dia tem demanda estrutural (instrumento errado: sobreaviso); "
                     "GA eventual usa sobreaviso corretamente. Do JSON (sobreaviso_por_ga, inteligencia.gas).")


def slide4(prs, dados, nar, ctx, blank):
    """Risco x custo-benefício: dimensionamento seguro por GA + a conta da conversão."""
    slide = prs.slides.add_slide(blank)
    s = nar["s4"]
    intel = dados.get("inteligencia") or {}
    chrome(slide, 4, 5, fmt(nar["fonte"], ctx))
    title_block(slide, fmt(s["eyebrow"], ctx), fmt(s["titulo"], ctx), title_size=27)

    # ---- esquerda: dimensionamento seguro por GA (plantão x demanda real)
    add_text(slide, "Dimensionamento seguro por GA — plantão hoje × demanda real de acionados",
             56, 152, 660, 18, size=13, color=NAVY, bold=True)
    add_text(slide, "Mediana e pico (p90/máx) de técnicos acionados no mesmo dia. A cobertura não cai: muda o instrumento.",
             56, 174, 660, 16, size=10, color=MUTED)
    gas_i = (intel.get("gas") or [])[:6]
    propostas = {
        "estrutural": ("Turno fixo + SA p/ pico", RED),
        "intermediario": ("SA enxuto (p90)", WARN),
        "eventual": ("Manter SA mínimo", GREEN_DARK),
    }
    table = add_table(slide, len(gas_i) + 1, 5, 56, 200, 660, 300, [190, 106, 118, 82, 164])
    cab = ["GA", "Plantão hoje", "Acionados/dia", "Pico", "Proposta"]
    for c, h in enumerate(cab):
        set_cell(table.cell(0, c), [(h, {"size": 10.5, "bold": True, "color": WHITE})],
                 fill=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    for r, g in enumerate(gas_i, start=1):
        rotulo, cor = propostas.get(g["classe"], ("—", MUTED))
        fill = WHITE if r % 2 else WASH
        set_cell(table.cell(r, 0), [(encurtar_nome(g["gestor"], 2), {"size": 10})], fill=fill)
        set_cell(table.cell(r, 1), [(f"{g['plantao_medio']:.0f} téc/dia",
                 {"size": 10, "align": PP_ALIGN.CENTER})], fill=fill)
        set_cell(table.cell(r, 2), [(f"med {g['acionados_mediana']} · p90 {g['acionados_p90']}",
                 {"size": 10, "align": PP_ALIGN.CENTER, "bold": True})], fill=fill)
        set_cell(table.cell(r, 3), [(f"{g['acionados_max']}",
                 {"size": 10, "align": PP_ALIGN.CENTER})], fill=fill)
        set_cell(table.cell(r, 4), [(rotulo, {"size": 9.5, "bold": True, "color": cor})], fill=fill)

    # régua de segurança
    add_rect(slide, 56, 516, 660, 62, RED_LIGHT, round_corners=True)
    add_text(slide, "REGRA DE PISO (RISCO)", 76, 526, 300, 14, size=10, color=RED, bold=True)
    add_text(slide, "Nenhuma GA reduz plantão abaixo do p90 de acionados simultâneos sem turno fixo no lugar. "
                    "O pico raro (máx) é coberto por backup combinado entre GAs vizinhas.",
             76, 542, 622, 32, size=10.5, color=INK, line_spacing=1.06)

    # ---- direita: a conta da conversão (número na coluna esquerda, rótulo na direita)
    ch = intel.get("custo_hora") or {}
    add_rect(slide, 740, 152, 482, 190, NAVY, round_corners=True)
    add_text(slide, "A CONTA DA CONVERSÃO", 762, 164, 400, 14, size=10.5, color="9FC9EE", bold=True)
    add_text(slide, f"R$ {ch.get('acionamento_rs_h', 0):.2f}".replace(".", ","),
             762, 186, 118, 30, size=22, bold=True, color=WHITE, font="Aptos Display")
    add_text(slide, "por hora acionada em sobreaviso (HE + adicionais)",
             890, 186, 312, 30, size=10.5, color="C9D6EF", anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, f"R$ {ch.get('turno_noturno_rs_h_estimado', 0):.2f}".replace(".", ","),
             762, 222, 118, 30, size=22, bold=True, color=GREEN, font="Aptos Display")
    add_text(slide, "por hora em turno fixo noturno (estimado*)",
             890, 222, 312, 30, size=10.5, color="C9D6EF", anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, fmt("Converter a demanda diária das GAs estruturais economiza "
                        "{economia_mes}/mês — sem perder cobertura.", ctx),
             762, 262, 440, 40, size=12.5, color=WHITE, bold=True, line_spacing=1.06)
    add_text(slide, "*hora-base derivada da própria espera (1/3 da hora normal) + 20% de adicional noturno. Validar com RH.",
             762, 312, 440, 24, size=8.5, color="8FA3CE")

    # ---- direita: nominal — quem já é turno de fato
    add_rect(slide, 740, 356, 482, 222, WHITE, line=LINE, round_corners=True)
    add_text(slide, fmt("Quem já trabalha o turno — só que pago como HE", ctx),
             762, 368, 440, 16, size=12, color=NAVY, bold=True)
    add_text(slide, fmt("Top 10 técnicos = {pct_top10} de todas as horas de acionamento. Os 5 maiores:", ctx),
             762, 388, 440, 16, size=10, color=MUTED)
    top5 = (intel.get("top_tecnicos") or [])[:5]
    ty = 408
    maxh_t = max((t["horas"] for t in top5), default=1) or 1
    for t in top5:
        add_text(slide, encurtar_nome(t["nome"], 2), 762, ty - 2, 150, 16, size=10, color=INK)
        add_rect(slide, 916, ty, 180, 12, "E7EDF5")
        add_rect(slide, 916, ty, max(180 * t["horas"] / maxh_t, 3), 12, BLUE_MID)
        add_text(slide, f"{t['horas']:.0f} h · {t['dias_acionado']} dias", 1104, ty - 2, 110, 16,
                 size=9.5, color=INK, bold=True)
        ty += 26
    add_text(slide, "Formalizar o turno primeiro para esse grupo: a demanda já é diária e comprovada.",
             762, ty + 2, 440, 18, size=10, color=MUTED)

    add_notes(slide, "SLIDE 4 — RISCO x CUSTO-BENEFÍCIO. Dimensionamento por GA (mediana=turno fixo, "
                     "p90=sobreaviso residual, máx=backup), regra de piso, conta da conversão e "
                     "nominal dos técnicos mais acionados. Do JSON (inteligencia).")


def slide5(prs, dados, nar, ctx, blank):
    slide = prs.slides.add_slide(blank)
    s = nar["s5"]
    chrome(slide, 5, 5, fmt(nar["fonte"], ctx))
    title_block(slide, fmt(s["eyebrow"], ctx), fmt(s["titulo"], ctx), title_size=27)

    # comparativo Orçado -> Proposta
    add_rect(slide, 56, 150, 360, 150, WHITE, line=LINE, round_corners=True)
    add_text(slide, "DO ORÇADO À PROPOSTA", 76, 162, 320, 14, size=10.5, color=BLUE, bold=True)
    add_text(slide, "Orçado do ciclo", 76, 188, 200, 16, size=11, color=MUTED)
    add_text(slide, ctx["orcado"], 76, 204, 320, 22, size=17, bold=True, color=INK)
    add_text(slide, "Proposta · −50%", 76, 234, 200, 16, size=11, color=MUTED)
    add_text(slide, ctx["meta"], 76, 250, 240, 28, size=21, bold=True, color=GREEN_DARK,
             font="Aptos Display")
    add_rect(slide, 300, 250, 96, 26, "E8F6EC", round_corners=True)
    add_text(slide, f"−{ctx['gap']}", 300, 250, 92, 26, size=11, bold=True, color=GREEN_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # alavancas
    ax = 436
    for al in s["alavancas"]:
        add_rect(slide, ax, 150, 256, 150, WASH, round_corners=True)
        add_text(slide, al["n"], ax + 18, 162, 60, 30, size=22, color=GREEN_DARK, bold=True, font="Aptos Display")
        add_text(slide, fmt(al["t"], ctx), ax + 18, 196, 224, 40, size=12, color=NAVY, bold=True, line_spacing=1.03)
        add_text(slide, fmt(al["d"], ctx), ax + 18, 236, 224, 58, size=9.5, color=INK, line_spacing=1.05)
        ax += 262

    # tabela FCA
    add_text(slide, "FCA — Fato · Causa · Ação", 56, 316, 600, 18, size=12.5, color=NAVY, bold=True)
    fca = s["fca"]
    table = add_table(slide, len(fca) + 1, 6, 56, 340, 1166, 300,
                      [250, 300, 316, 110, 150, 40])
    cab = ["Fato", "Causa", "Ação", "Prazo", "Responsável", ""]
    for c, h in enumerate(cab):
        set_cell(table.cell(0, c), [(h, {"size": 10.5, "bold": True, "color": WHITE})],
                 fill=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    cores_status = {"urgent": (RED, "Imediato"), "start": (BLUE, "A iniciar"),
                    "watch": (WARN, "Atenção"), "validate": (AMBER, "Validar")}
    for r, item in enumerate(fca, start=1):
        fill = WHITE if r % 2 else WASH
        set_cell(table.cell(r, 0), [(fmt(item["fato"], ctx), {"size": 10, "bold": True, "color": NAVY})], fill=fill, anchor=MSO_ANCHOR.TOP)
        set_cell(table.cell(r, 1), [(fmt(item["causa"], ctx), {"size": 9.5})], fill=fill, anchor=MSO_ANCHOR.TOP)
        set_cell(table.cell(r, 2), [(fmt(item["acao"], ctx), {"size": 9.5})], fill=fill, anchor=MSO_ANCHOR.TOP)
        set_cell(table.cell(r, 3), [(fmt(item["prazo"], ctx), {"size": 9.5, "bold": True})], fill=fill, anchor=MSO_ANCHOR.TOP)
        set_cell(table.cell(r, 4), [(fmt(item["resp"], ctx), {"size": 9.5})], fill=fill, anchor=MSO_ANCHOR.TOP)
        cor, rotulo = cores_status.get(item.get("status", "start"), (BLUE, "A iniciar"))
        set_cell(table.cell(r, 5), [("●", {"size": 12, "color": cor, "align": PP_ALIGN.CENTER})],
                 fill=fill, anchor=MSO_ANCHOR.MIDDLE)

    add_notes(slide, "SLIDE 5 — AÇÃO (proposta + FCA cirúrgico). Orçado->Proposta (-50%), 3 alavancas "
                     "e a tabela FCA com ações específicas por GA. Edite textos em narrativa_<regiao>.json (s5).")


# ------------------------------------------------------------------- contexto
def montar_contexto(dados, nar):
    m = dados["meta"]
    ro = dados["real_orcado_mes"]
    cesta = dados["cesta"]
    total_cesta = cesta_get(cesta, "Total") or sum(
        v for k, v in cesta.items() if str(k).lower() != "total")
    sobreaviso = (cesta_get(cesta, "ACIONAMENTO SOBRE AVISO") or 0) + (cesta_get(cesta, "ESPERA SOBRE AVISO") or 0)
    dsr = cesta_get(cesta, "DSR") or 0
    real = ro["real"]
    orcado = m["orcado_mes"]
    meta_val = m["meta_mes"]
    gap = real - meta_val
    o = dados["ociosidade"]
    ga = dados["sobreaviso_por_ga"]
    tot_ga = sum(g["horas_total"] for g in ga) or 1
    top3 = sum(g["horas_total"] for g in ga[:3]) / tot_ga

    ctx = {
        "regiao": nar["regiao"], "codigo": nar["codigo"],
        "mes_foto": nar["mes_foto"], "data_foto": nar["data_foto"],
        "orcado": moeda(orcado), "meta": moeda(meta_val), "real": moeda(real),
        "meta_val": meta_val,
        "gap": moeda(gap), "gap_pct": pct(gap / real if real else 0),
        "var_orcado": pct(abs((real - orcado) / orcado) if orcado else 0) + (" abaixo" if real < orcado else " acima"),
        "sobreaviso": moeda(sobreaviso), "sobreaviso_pct": pct(sobreaviso / total_cesta if total_cesta else 0),
        "dsr_pct": pct(dsr / total_cesta if total_cesta else 0),
        "ocio_hc_espera": f"{o['hc_espera_medio']:.1f}".replace(".", ","),
        "ocio_hc_acion": f"{o['hc_acionado_medio']:.1f}".replace(".", ","),
        "ocio_pct_hc": pct(o.get("pct_acionamento_hc")),
        "ocio_horas_espera": f"{o['horas_espera_total']:,.0f} h".replace(",", "."),
        "ocio_horas_acion": f"{o['horas_acionamento_total']:,.0f} h".replace(",", "."),
        "ocio_pct_horas": pct(o.get("pct_horas_acionamento")),
        "top3_ga_pct": pct(top3),
    }

    # faixa horária de maior/menor acionamento (data-driven, já ordenada desc)
    faixas = dados.get("faixa_horaria") or []
    if faixas:
        tot_f = sum(f["horas"] for f in faixas) or 1
        ctx["faixa_top"] = faixas[0]["faixa"]
        ctx["faixa_top_pct"] = pct(faixas[0]["horas"] / tot_f)
        ctx["faixa_baixas"] = " e ".join(f["faixa"].lower() for f in faixas[-2:])
    else:
        ctx["faixa_top"] = ctx["faixa_top_pct"] = ctx["faixa_baixas"] = "—"

    # inteligência: regimes, GAs estruturais, conversão p/ turno fixo
    intel = dados.get("inteligencia") or {}
    reg = intel.get("regimes") or {}
    ch = intel.get("custo_hora") or {}
    gas_i = intel.get("gas") or []
    dias = intel.get("dias_periodo") or 30
    estruturais = [g for g in gas_i if g["classe"] == "estrutural"]
    eventuais = [g for g in gas_i if g["classe"] == "eventual"]
    ctx["pct_fds"] = pct(reg.get("pct_horas_fds"))
    ctx["h_util_noturno"] = f"{reg.get('horas_util_noturno', 0):,.0f} h".replace(",", ".")
    ctx["h_fds_diurno"] = f"{reg.get('horas_fds_diurno', 0):,.0f} h".replace(",", ".")
    ctx["n_estruturais"] = str(len(estruturais))
    ctx["gas_estruturais"] = ", ".join(encurtar_nome(g["gestor"]) for g in estruturais[:3]) or "—"
    ctx["gas_eventuais"] = ", ".join(encurtar_nome(g["gestor"]) for g in eventuais[:3]) or "—"
    ctx["pct_top10"] = pct(intel.get("pct_horas_top10"))
    ctx["rs_acion_h"] = f"R$ {ch.get('acionamento_rs_h', 0):.2f}".replace(".", ",")
    ctx["rs_turno_h"] = f"R$ {ch.get('turno_noturno_rs_h_estimado', 0):.2f}".replace(".", ",")
    ctx["pico_max"] = str(max((g["acionados_max"] for g in gas_i), default=0))
    # economia estimada/mês da conversão nas GAs estruturais:
    # horas acionadas × (R$/h acionamento − R$/h turno) + metade da espera eliminada
    delta = (ch.get("acionamento_rs_h") or 0) - (ch.get("turno_noturno_rs_h_estimado") or 0)
    econ_periodo = sum(g["horas_acionamento"] for g in estruturais) * max(delta, 0) \
        + 0.5 * sum(g["custo_espera"] for g in estruturais)
    ctx["economia_mes"] = moeda(econ_periodo / dias * 30)
    ctx["custo_estruturais_mes"] = moeda(sum(g["custo_total"] for g in estruturais) / dias * 30)
    return ctx


def main():
    ap = argparse.ArgumentParser(description="Monta o PPTX de Custo de Horas (Alloha).")
    ap.add_argument("--regiao", required=True, help="Slug, ex.: ceara. Usa dados_<slug>.json e narrativa_<slug>.json.")
    ap.add_argument("--dados", default=None)
    ap.add_argument("--narrativa", default=None)
    ap.add_argument("--saida", default=None)
    args = ap.parse_args()

    dados_path = Path(args.dados) if args.dados else ROOT / "outputs" / "horas" / f"dados_{args.regiao}.json"
    nar_path = Path(args.narrativa) if args.narrativa else ROOT / "docs" / "horas" / f"narrativa_{args.regiao}.json"
    dados = json.loads(dados_path.read_text(encoding="utf-8"))
    nar = json.loads(nar_path.read_text(encoding="utf-8"))
    ctx = montar_contexto(dados, nar)

    prs = Presentation()
    prs.slide_width = px(1280)
    prs.slide_height = px(720)
    blank = prs.slide_layouts[6]
    slide1(prs, dados, nar, ctx, blank)
    slide2(prs, dados, nar, ctx, blank)
    slide3(prs, dados, nar, ctx, blank)
    slide4(prs, dados, nar, ctx, blank)
    slide5(prs, dados, nar, ctx, blank)

    slug = "".join(c for c in unicodedata.normalize("NFKD", args.regiao)
                   if not unicodedata.combining(c)).replace(" ", "_")
    saida = Path(args.saida) if args.saida else ROOT / "outputs" / "horas" / f"Horas_{slug.capitalize()}_Executivo.pptx"
    saida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(saida))
    print(f"OK -> {saida}")


if __name__ == "__main__":
    main()
