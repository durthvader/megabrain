"""Gera outputs/3ps/3Ps_MaranhaoPiaui_Executivo.pptx com a identidade visual Alloha.

Espelha a versão final revisada manualmente (3Ps_MaranhaoPiaui_Executivo (VF).pptx),
sem comparações com o Ceará -- cada painel regional é autônomo.
Uso: python scripts/build_3ps_presentation_mapi.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "outputs" / "3ps" / "3Ps_MaranhaoPiaui_Executivo.pptx"

NAVY = "081551"
BLUE = "1476C9"
BLUE_LIGHT = "DCECF9"
GREEN = "42DF4B"
GREEN_DARK = "13853A"
RED = "D94545"
RED_LIGHT = "FDE2DF"
ORANGE = "F08A24"
AMBER = "E4BD30"
INK = "182238"
MUTED = "667085"
LINE = "D9E0EA"
WASH = "F3F6FA"
WHITE = "FFFFFF"


def px(value: float) -> Emu:
    return Emu(int(value * 9525))


def rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


def add_rect(slide, left, top, width, height, fill, line=None, line_width=1.0, round_corners=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_corners else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, px(left), px(top), px(width), px(height))
    if round_corners:
        shape.adjustments[0] = 0.07
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_text(slide, value, left, top, width, height, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Aptos", line_spacing=1.06):
    box = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, margin, 0)
    lines = value if isinstance(value, list) else [(value, {})]
    for index, (text, opts) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = opts.get("align", align)
        paragraph.line_spacing = opts.get("line_spacing", line_spacing)
        if index > 0:
            paragraph.space_before = Pt(opts.get("space_before", 2))
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        run.font.color.rgb = rgb(opts.get("color", color))
        run.font.name = opts.get("font", font)
    return box


def chrome(slide, number, total, source):
    add_rect(slide, 0, 0, 1100, 8, NAVY)
    add_rect(slide, 1100, 0, 180, 8, GREEN)
    add_text(slide, "alloha", 1090, 30, 130, 36, size=27, color=NAVY,
             align=PP_ALIGN.RIGHT, font="Aptos Display")
    add_text(slide, "F I B R A", 1090, 66, 130, 14, size=8, color=GREEN_DARK,
             bold=True, align=PP_ALIGN.RIGHT)
    if source:
        add_text(slide, source, 56, 688, 1040, 16, size=9, color="8B94A6")
    add_text(slide, f"{number} / {total}", 1150, 688, 74, 16, size=9,
             color="8B94A6", align=PP_ALIGN.RIGHT)


def title_block(slide, eyebrow, title, lead=None, title_size=34, lead_top=138):
    add_text(slide, eyebrow, 56, 36, 900, 20, size=12, color=BLUE, bold=True)
    add_text(slide, title, 56, 60, 1000, 76, size=title_size, color=NAVY,
             bold=True, font="Aptos Display", line_spacing=1.0)
    if lead:
        add_text(slide, lead, 56, lead_top, 1120, 44, size=15, color="45516A")


def set_cell(cell, lines, fill=None, anchor=MSO_ANCHOR.MIDDLE, v_margin=6):
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill)
    cell.vertical_anchor = anchor
    cell.margin_left = px(10)
    cell.margin_right = px(8)
    cell.margin_top = px(v_margin)
    cell.margin_bottom = px(v_margin)
    tf = cell.text_frame
    tf.word_wrap = True
    for index, (text, opts) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = opts.get("align", PP_ALIGN.LEFT)
        paragraph.line_spacing = opts.get("line_spacing", 1.05)
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", 11))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = rgb(opts.get("color", INK))
        run.font.name = opts.get("font", "Aptos")


def add_table(slide, rows, cols, left, top, width, height, col_widths):
    frame = slide.shapes.add_table(rows, cols, px(left), px(top), px(width), px(height))
    table = frame.table
    table.first_row = False
    table.horz_banding = False
    for index, cw in enumerate(col_widths):
        table.columns[index].width = px(cw)
    return table


def hbar(slide, label, value_text, left, top, label_width, bar_width, ratio,
         color=BLUE, track="E7EDF5", bar_height=12):
    add_text(slide, label, left, top - 2, label_width, 16, size=10.5, color=INK)
    track_left = left + label_width + 8
    add_rect(slide, track_left, top, bar_width, bar_height, track)
    add_rect(slide, track_left, top, max(bar_width * ratio, 4), bar_height, color)
    add_text(slide, value_text, track_left + bar_width + 8, top - 2, 44, 16,
             size=11, color=INK, bold=True)


prs = Presentation()
prs.slide_width = px(1280)
prs.slide_height = px(720)
blank = prs.slide_layouts[6]

# ---------------------------------------------------------------- Slide 1
slide = prs.slides.add_slide(blank)
chrome(slide, 1, 4, "Fonte: exportações de BI operacional da regional R7.2 (MA/PI), 01–09/07/2026 · OS = ordem de serviço concluída com baixa OK")
title_block(
    slide,
    "OPERAÇÕES · MARANHÃO/PIAUÍ (R7.2) · 09/07/2026",
    "Maranhão volta a cumprir o Prazo 24h; Produção segue abaixo do ideal",
    "No fechamento de 08/07, o Maranhão (C.32) passou a cumprir a meta de Prazo 24h, mas segue abaixo do esperado "
    "em Produção; o Piauí (C.31) está com Presença, Produção e Prazo na meta.",
    title_size=28,
    lead_top=152,
)

table = add_table(slide, 3, 5, 56, 202, 1168, 204, [200, 236, 250, 236, 246])
headers = [
    ("Cluster", ""),
    ("Presença", "meta mínima 95%"),
    ("Produção", "% da meta esperada, 08/07"),
    ("Prazo 24h", "meta 85% das ordens"),
    ("Fila acumulada", "em dias de trabalho, 09/07"),
]
for col, (head, sub) in enumerate(headers):
    lines = [(head, {"size": 13, "bold": True, "color": WHITE})]
    if sub:
        lines.append((sub, {"size": 9.5, "color": "BFC9E3"}))
    set_cell(table.cell(0, col), lines, fill=NAVY)

row_c31 = [
    [("C.31", {"size": 16, "bold": True, "color": NAVY}), ("Piauí", {"size": 10, "color": MUTED})],
    [("114,3%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("19,3 p.p. acima da meta", {"size": 10, "color": MUTED})],
    [("144,8%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("acima do esperado", {"size": 10, "color": MUTED})],
    [("92,0%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("7,0 p.p. acima da meta", {"size": 10, "color": MUTED})],
    [("0,74", {"size": 16, "bold": True, "color": GREEN_DARK}), ("dentro da meta", {"size": 10, "color": MUTED})],
]
row_c32 = [
    [("C.32", {"size": 16, "bold": True, "color": NAVY}), ("Maranhão", {"size": 10, "color": MUTED})],
    [("100,0%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("5,0 p.p. acima da meta", {"size": 10, "color": MUTED})],
    [("88,4%", {"size": 16, "bold": True, "color": RED}), ("abaixo do esperado", {"size": 10, "color": MUTED})],
    [("89,4%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("4,4 p.p. acima da meta", {"size": 10, "color": MUTED})],
    [("0,72", {"size": 16, "bold": True, "color": GREEN_DARK}), ("dentro da meta", {"size": 10, "color": MUTED})],
]
for col, lines in enumerate(row_c31):
    set_cell(table.cell(1, col), lines, fill=WHITE)
for col, lines in enumerate(row_c32):
    set_cell(table.cell(2, col), lines, fill="EAF2FB")

add_rect(slide, 56, 470, 545, 168, NAVY, round_corners=True)
add_text(slide, "DECISÃO OPERACIONAL", 80, 484, 320, 18, size=11, color="9FC9EE", bold=True)
add_text(
    slide,
    "Manter o foco no Maranhão: Produção segue abaixo do esperado. Execução também nominal: 1 GA concentra o "
    "maior gap; 9 técnicos exigem decisão individual (8 com maior desvio de produção e 1 sem nenhuma baixa).",
    80, 506, 498, 118, size=14, color=WHITE, bold=True,
)

kpis = [
    ("18 de 68", "técnicos estão em Q3/Q4 de produtividade (1:1 OK) — abaixo da mediana ou sem produção"),
    ("-7,1%", "queda de Produção 1:1 OK entre a semana fechada e a semana em curso"),
]
for index, (value, label) in enumerate(kpis):
    left = 625 + index * 320
    add_rect(slide, left, 470, 290, 168, WHITE, line=LINE, round_corners=True)
    add_text(slide, value, left + 18, 486, 254, 32, size=23, color=BLUE, bold=True)
    add_text(slide, label, left + 18, 522, 254, 104, size=11, color=MUTED)

# ---------------------------------------------------------------- Slide 2
slide = prs.slides.add_slide(blank)
chrome(slide, 2, 4, "Fonte: Analítico Nominal 01–09/07/2026, Quartil de Produtividade S27/S28 e exportações de motivo/submotivo de improdutividade da R7.2")
title_block(
    slide,
    "DIAGNÓSTICO DE PRODUÇÃO",
    "A produtividade está concentrada na faixa mediana — o problema é a cauda, não o time inteiro",
    title_size=27,
)

add_rect(slide, 56, 150, 610, 372, WHITE, line=LINE, round_corners=True)
add_text(slide, "68 técnicos elegíveis", 78, 166, 330, 20, size=15, color=NAVY, bold=True)
add_text(slide, "Perfil quase todo veterano", 340, 168, 306, 18, size=12, color=MUTED, align=PP_ALIGN.RIGHT)
add_text(
    slide,
    "Quartil de produtividade (1:1 OK), média das semanas S27 e S28. Quadro concentrado em >90 dias de casa: "
    "16 de 18 no Piauí, 46 de 50 no Maranhão.",
    78, 190, 566, 28, size=10, color=MUTED,
)
segments = [(35.6, GREEN_DARK), (31.1, AMBER), (28.8, ORANGE), (4.5, RED)]
bar_left, bar_width = 78, 566
offset = bar_left
for pct, color in segments:
    seg_width = bar_width * pct / 100
    add_rect(slide, offset, 222, seg_width, 20, color)
    offset += seg_width
legend = [
    ("35,6%", "Q1 — acima da mediana", GREEN_DARK),
    ("31,1%", "Q2 — mediana", AMBER),
    ("28,8%", "Q3 — abaixo da mediana", ORANGE),
    ("4,5%", "Q4 — baixa ou sem produção", RED),
]
for index, (count, label, color) in enumerate(legend):
    left = 78 + index * 144
    add_rect(slide, left, 252, 7, 28, color)
    add_text(slide, [(count, {"size": 14, "bold": True}), (label, {"size": 8.5, "color": MUTED})],
             left + 12, 250, 132, 34)
add_rect(slide, 78, 292, 566, 66, "FDF4EA")
add_rect(slide, 78, 292, 5, 66, ORANGE)
add_text(
    slide,
    "A maior parte do quadro está em Q1/Q2. O risco não é generalizado: está concentrado em poucos nomes e em "
    "uma área específica.",
    92, 300, 544, 56, size=10.5, color="5A4632",
)

add_rect(slide, 686, 150, 538, 372, WHITE, line=LINE, round_corners=True)
add_text(slide, "273 visitas improdutivas", 708, 166, 320, 20, size=15, color=NAVY, bold=True)
add_text(slide, "174 com motivo registrado", 866, 168, 338, 18, size=12, color=RED, bold=True, align=PP_ALIGN.RIGHT)
add_text(
    slide,
    "Motivos registrados nos 4 maiores grupos (01–09/07):",
    708, 190, 494, 30, size=10, color=MUTED,
)
causes = [
    ("Abertura indevida", "39", 1.0),
    ("Falha massiva", "37", 0.949),
    ("Cliente ausente", "20", 0.513),
    ("Solicitação de reagendamento", "17", 0.436),
]
for index, (label, value, ratio) in enumerate(causes):
    hbar(slide, label, value, 708, 236 + index * 34, 190, 240, ratio, color=ORANGE)
add_rect(slide, 708, 380, 494, 78, "F1F9F1")
add_rect(slide, 708, 380, 5, 78, GREEN)
add_text(slide, "74%", 724, 390, 90, 30, size=19, color=GREEN_DARK, bold=True)
add_text(
    slide,
    "das improdutivas classificadas estão nos 5 principais motivos, com “Abertura indevida” e “Falha massiva”, "
    "parte relevante nasce antes do despacho.",
    822, 388, 368, 60, size=10.5, color="4A5B52",
)
add_text(
    slide,
    "Como ler: Abertura indevida são os casos de Serviço não é mais necessário (19), Cliente não solicitou "
    "serviço (12), Demanda não é técnica (5) e Serviço não autorizado (3).",
    708, 466, 494, 44, size=8.5, color=MUTED,
)
solutions = [
    ("01", "Reverter a queda S27→S28",
     "Produção 1:1 OK caiu de 3,25 para 3,02 OS/dia. Acompanhamento diário para retomar o patamar da semana "
     "anterior até o fim da S28."),
    ("02", "Plano individual nominal",
     "Foco nos 8 maiores gaps individuais e no GA com maior concentração de perda. Detalhe na página seguinte."),
    ("03", "Qualidade da demanda",
     "Travar abertura de OS sem validação técnica e tratar falha massiva antes do despacho — os dois motivos "
     "somam 76 das 174 baixas classificadas."),
]
for index, (number, head, body) in enumerate(solutions):
    left = 56 + index * 398
    add_rect(slide, left, 540, 382, 118, WASH)
    add_rect(slide, left, 540, 382, 3, BLUE)
    add_text(slide, number, left + 14, 552, 40, 26, size=20, color="A2B2C6", bold=True)
    add_text(slide, head, left + 54, 552, 316, 20, size=13, color=NAVY, bold=True)
    add_text(slide, body, left + 54, 574, 316, 78, size=10, color=MUTED)

# ---------------------------------------------------------------- Slide 3
slide = prs.slides.add_slide(blank)
chrome(slide, 3, 4, "Fonte: Analítico Nominal por técnico, 01 a 09/07/2026 · 68 técnicos elegíveis · gap estimado, dias trabalhados inferidos")
title_block(
    slide,
    "GESTÃO NOMINAL · 01 A 09/07/2026",
    "Uma área e oito técnicos concentram a maior parte do gap do Maranhão",
    title_size=27,
)

add_rect(slide, 56, 148, 452, 400, WHITE, line=LINE, round_corners=True)
add_text(slide, "Gap por área de gestão (GA)", 78, 164, 300, 20, size=14, color=NAVY, bold=True)
add_text(slide, "em OS não executadas (estimado)", 300, 166, 188, 18, size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)
gas = [
    ("Naan Carlos Oliveira · C.32", "95", 1.0, BLUE),
    ("Leonardo Silva · C.31/C.32", "83", 0.874, BLUE),
    ("Francisco Sales Jr. · C.32", "80", 0.842, BLUE),
    ("Roberval Silva · C.32", "76", 0.8, "78A9D8"),
    ("Ruan Labre da Silva · C.32", "49", 0.516, "78A9D8"),
]
for index, (label, value, ratio, color) in enumerate(gas):
    hbar(slide, label, value, 78, 202 + index * 32, 172, 180, ratio, color=color)
add_rect(slide, 78, 366, 408, 96, "F1F9F1")
add_rect(slide, 78, 366, 5, 96, GREEN)
add_text(slide, "1 técnico", 94, 378, 120, 30, size=18, color=GREEN_DARK, bold=True)
add_text(
    slide,
    "ficou sem nenhuma baixa no período inteiro (01–09/07): Israel Costa da Silva (C.32, GA Naan Carlos Cabral "
    "Oliveira).",
    94, 410, 376, 44, size=10.5, color="4A5B52",
)

table = add_table(slide, 3, 3, 528, 148, 696, 264, [232, 240, 224])
set_cell(table.cell(0, 0), [("TÉCNICOS", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
set_cell(table.cell(0, 1), [("O QUE OS DADOS MOSTRAM", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
set_cell(table.cell(0, 2), [("AÇÃO IMEDIATA", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
nominal_rows = [
    (
        [("1 técnico sem baixa no período", {"size": 10, "bold": True, "color": NAVY}),
         ("Israel Costa da Silva (C.32)", {"size": 8, "color": MUTED})],
        [("Sem escada para atividade.", {"size": 8.5})],
        [("Escada enviada para Grajaú-MA, previsão para 13/07.", {"size": 8.5})],
        WHITE,
    ),
    (
        [("8 maiores desvios individuais", {"size": 10, "bold": True, "color": NAVY}),
         ("Andre Santos Pereira e Natan Padilha Pinheiro (gap 19); David Wesllem de Sousa Baldez (17); Marcio Andre Cutrim Costa e Lucas dos Santos Cutrim (14); Antonio Pedro de Oliveira, Gabriel Leite Barbosa e Felipe Sousa Quirino (13) — todos C.32", {"size": 8, "color": MUTED})],
        [("Concentram boa parte do gap do Maranhão; a maioria com poucos dias trabalhados no período (4 a 7 de 9).", {"size": 8.5})],
        [("Conversa individual gestor-técnico esta semana, com plano simples de recuperação e verificação da escala real.", {"size": 8.5})],
        "F8FAFD",
    ),
]
for row_index, (col_a, col_b, col_c, fill) in enumerate(nominal_rows, start=1):
    set_cell(table.cell(row_index, 0), col_a, fill=fill, v_margin=4)
    set_cell(table.cell(row_index, 1), col_b, fill=fill, v_margin=4)
    set_cell(table.cell(row_index, 2), col_c, fill=fill, v_margin=4)
for row in table.rows:
    row.height = px(88)

tiers = [
    ("1", "sem baixa no período — resolver em 24h", RED),
    ("8", "desvio crítico, 29% da perda — plano individual já", RED),
    ("22", "produtividade baixa (< 70% da meta) — conversa até 12/07", ORANGE),
    ("23", "produtividade média (70 a 100% da meta) — ajuste de carga e rota", AMBER),
    ("14", "na meta (100% ou mais) — reconhecer e usar de referência", GREEN_DARK),
]
for index, (count, label, color) in enumerate(tiers):
    left = 56 + index * 238
    add_rect(slide, left, 588, 226, 78, WASH)
    add_rect(slide, left, 588, 226, 3, color)
    add_text(slide, count, left + 14, 602, 52, 30, size=22, color=color, bold=True)
    add_text(slide, label, left + 70, 598, 148, 60, size=9.5, color=MUTED)

# ---------------------------------------------------------------- Slide 4
slide = prs.slides.add_slide(blank)
chrome(slide, 4, 4, "Fonte: Analítico Nominal 01–09/07/2026 e apontamento operacional da regional R7.2")
title_block(
    slide,
    "FCA · FATO, CAUSA E AÇÃO",
    "Duas frentes colocam o Maranhão em recuperação da produção",
    title_size=27,
)

table = add_table(slide, 3, 6, 42, 148, 1196, 280, [200, 258, 268, 116, 216, 138])
fca_heads = ["FATO", "CAUSA", "AÇÃO", "PRAZO", "RESPONSÁVEL", "STATUS"]
for col, head in enumerate(fca_heads):
    set_cell(table.cell(0, col), [(head, {"size": 10.5, "bold": True, "color": WHITE})], fill=NAVY)

fca_rows = [
    (
        [("Produção abaixo da meta", {"size": 11, "bold": True, "color": NAVY}),
         ("Desvio em 11,6pp", {"size": 9, "color": MUTED})],
        [("Um técnico sem escada em cidade Grajaú-MA.", {"size": 9.5})],
        [("Fornecimento de escada. CD de Teresina.", {"size": 9.5})],
        [("14/07", {"size": 11, "bold": True, "color": NAVY})],
        [("Naan Carlos (GA)", {"size": 9.5})],
        ("Em Andamento", BLUE_LIGHT, "114F88"),
        WHITE,
    ),
    (
        [("Gap concentrado em 8 técnicos", {"size": 11, "bold": True, "color": NAVY}),
         ("122 OS de desvio (29% do gap total)", {"size": 9, "color": MUTED})],
        [("8 técnicos com maior desvio individual devido a baixa demanda nas cidades de atuação e falta de materiais para preventivas.", {"size": 9.5})],
        [("Complementar rota dos técnicos com preventivas, clean up e presets.", {"size": 9.5})],
        [("16/07", {"size": 11, "bold": True, "color": NAVY})],
        [("Wellington Figueira (GG)", {"size": 9.5})],
        ("Em Andamento", BLUE_LIGHT, "114F88"),
        "F8FAFD",
    ),
]
for row_index, (fato, causa, acao, prazo, resp, status, fill) in enumerate(fca_rows, start=1):
    set_cell(table.cell(row_index, 0), fato, fill=fill)
    set_cell(table.cell(row_index, 1), causa, fill=fill)
    set_cell(table.cell(row_index, 2), acao, fill=fill)
    set_cell(table.cell(row_index, 3), prazo, fill=fill)
    set_cell(table.cell(row_index, 4), resp, fill=fill)
    status_text, status_fill, status_color = status
    set_cell(
        table.cell(row_index, 5),
        [(status_text, {"size": 10, "bold": True, "color": status_color, "align": PP_ALIGN.CENTER})],
        fill=status_fill,
    )
for row in table.rows:
    row.height = px(90)

add_rect(slide, 42, 460, 1196, 68, NAVY)
add_rect(slide, 42, 460, 7, 68, GREEN)
add_text(slide, "DECISÕES SOLICITADAS", 64, 476, 200, 18, size=11.5, color=WHITE, bold=True)
add_text(
    slide,
    "Confirmar entrega da escada em Grajaú-MA até 14/07   ·   Aprovar plano de rota complementar para os 8 "
    "técnicos até 16/07   ·   Validar responsáveis e prazos",
    64, 496, 1146, 28, size=10.5, color="D5DEF3",
)

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(f"OK: {OUTPUT}")
