"""Gera outputs/3ps/3Ps_Ceara_Executivo.pptx com a identidade visual Alloha.

Substitui o antigo build_3ps_presentation.mjs (dependia do runtime do Codex).
Uso: python scripts/build_3ps_presentation.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "outputs" / "3ps" / "3Ps_Ceara_Executivo.pptx"

NAVY = "081551"
BLUE = "1476C9"
BLUE_LIGHT = "DCECF9"
GREEN = "42DF4B"
GREEN_DARK = "13853A"
RED = "D94545"
RED_LIGHT = "FDE2DF"
ORANGE = "F08A24"
AMBER = "E4BD30"
INK = "182238"
MUTED = "667085"
LINE = "D9E0EA"
WASH = "F3F6FA"
WHITE = "FFFFFF"


def px(value: float) -> Emu:
    return Emu(int(value * 9525))


def rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


def add_rect(slide, left, top, width, height, fill, line=None, line_width=1.0, round_corners=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_corners else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, px(left), px(top), px(width), px(height))
    if round_corners:
        shape.adjustments[0] = 0.07
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_text(slide, value, left, top, width, height, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Aptos", line_spacing=1.06):
    """value: str ou lista de (texto, opções por linha)."""
    box = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, margin, 0)
    lines = value if isinstance(value, list) else [(value, {})]
    for index, (text, opts) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = opts.get("align", align)
        paragraph.line_spacing = opts.get("line_spacing", line_spacing)
        if index > 0:
            paragraph.space_before = Pt(opts.get("space_before", 2))
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        run.font.color.rgb = rgb(opts.get("color", color))
        run.font.name = opts.get("font", font)
    return box


def chrome(slide, number, total, source):
    add_rect(slide, 0, 0, 1100, 8, NAVY)
    add_rect(slide, 1100, 0, 180, 8, GREEN)
    add_text(slide, "alloha", 1090, 30, 130, 36, size=27, color=NAVY,
             align=PP_ALIGN.RIGHT, font="Aptos Display")
    add_text(slide, "F I B R A", 1090, 66, 130, 14, size=8, color=GREEN_DARK,
             bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, source, 56, 688, 1040, 16, size=9, color="8B94A6")
    add_text(slide, f"{number} / {total}", 1150, 688, 74, 16, size=9,
             color="8B94A6", align=PP_ALIGN.RIGHT)


def title_block(slide, eyebrow, title, lead=None, title_size=34, lead_top=138):
    add_text(slide, eyebrow, 56, 36, 900, 20, size=12, color=BLUE, bold=True)
    add_text(slide, title, 56, 60, 1000, 76, size=title_size, color=NAVY,
             bold=True, font="Aptos Display", line_spacing=1.0)
    if lead:
        add_text(slide, lead, 56, lead_top, 1120, 44, size=15, color="45516A")


def set_cell(cell, lines, fill=None, anchor=MSO_ANCHOR.MIDDLE, v_margin=6):
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill)
    cell.vertical_anchor = anchor
    cell.margin_left = px(10)
    cell.margin_right = px(8)
    cell.margin_top = px(v_margin)
    cell.margin_bottom = px(v_margin)
    tf = cell.text_frame
    tf.word_wrap = True
    for index, (text, opts) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = opts.get("align", PP_ALIGN.LEFT)
        paragraph.line_spacing = opts.get("line_spacing", 1.05)
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", 11))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = rgb(opts.get("color", INK))
        run.font.name = opts.get("font", "Aptos")


def add_table(slide, rows, cols, left, top, width, height, col_widths):
    frame = slide.shapes.add_table(rows, cols, px(left), px(top), px(width), px(height))
    table = frame.table
    table.first_row = False
    table.horz_banding = False
    for index, cw in enumerate(col_widths):
        table.columns[index].width = px(cw)
    return table


def hbar(slide, label, value_text, left, top, label_width, bar_width, ratio,
         color=BLUE, track="E7EDF5", bar_height=12):
    add_text(slide, label, left, top - 2, label_width, 16, size=10.5, color=INK)
    track_left = left + label_width + 8
    add_rect(slide, track_left, top, bar_width, bar_height, track)
    add_rect(slide, track_left, top, max(bar_width * ratio, 4), bar_height, color)
    add_text(slide, value_text, track_left + bar_width + 8, top - 2, 44, 16,
             size=11, color=INK, bold=True)


prs = Presentation()
prs.slide_width = px(1280)
prs.slide_height = px(720)
blank = prs.slide_layouts[6]

# ---------------------------------------------------------------- Slide 1
slide = prs.slides.add_slide(blank)
chrome(slide, 1, 4, "Fonte: BI operacional e bases nominais do projeto 3Ps · OS = ordem de serviço concluída com baixa OK")
title_block(
    slide,
    "OPERAÇÕES · CEARÁ · 06/07/2026",
    "Fortaleza concentra o risco; o Interior precisa proteger o Prazo",
    "Fortaleza (C.27) está abaixo da meta nos três indicadores. O Interior (C.26) cumpre Presença e fechou o Prazo "
    "no limite, mas começou a semana com fila de ordens acima de um dia de trabalho.",
    title_size=32,
    lead_top=152,
)

table = add_table(slide, 3, 5, 56, 202, 1168, 204, [210, 226, 250, 236, 246])
headers = [
    ("Cluster", ""),
    ("Presença", "meta mínima 95%"),
    ("Produção OK", "semana 29/06–05/07"),
    ("Prazo 24h", "meta 85% das ordens"),
    ("Fila acumulada", "em dias de trabalho, 06/07"),
]
for col, (head, sub) in enumerate(headers):
    lines = [(head, {"size": 13, "bold": True, "color": WHITE})]
    if sub:
        lines.append((sub, {"size": 9.5, "color": "BFC9E3"}))
    set_cell(table.cell(0, col), lines, fill=NAVY)

row_c26 = [
    [("C.26", {"size": 16, "bold": True, "color": NAVY}), ("Ceará Interior", {"size": 10, "color": MUTED})],
    [("102,5%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("acima da meta", {"size": 10, "color": MUTED})],
    [("3,5 OS/dia", {"size": 16, "bold": True, "color": "9A6200"}), ("veteranos abaixo da meta de 4", {"size": 10, "color": MUTED})],
    [("85,2%", {"size": 16, "bold": True, "color": "9A6200"}), ("só 0,2 p.p. acima da meta", {"size": 10, "color": MUTED})],
    [("1,125 dia", {"size": 16, "bold": True, "color": RED}), ("acima do limite de 1 dia", {"size": 10, "color": MUTED})],
]
row_c27 = [
    [("C.27", {"size": 16, "bold": True, "color": NAVY}), ("Fortaleza", {"size": 10, "color": MUTED})],
    [("94,5%", {"size": 16, "bold": True, "color": RED}), ("0,5 p.p. abaixo da meta", {"size": 10, "color": MUTED})],
    [("3,2 OS/dia", {"size": 16, "bold": True, "color": RED}), ("61 dos 81 técnicos abaixo", {"size": 10, "color": MUTED})],
    [("83,2%", {"size": 16, "bold": True, "color": RED}), ("1,8 p.p. abaixo da meta", {"size": 10, "color": MUTED})],
    [("0,820 dia", {"size": 16, "bold": True, "color": GREEN_DARK}), ("dentro do limite", {"size": 10, "color": MUTED})],
]
for col, lines in enumerate(row_c26):
    set_cell(table.cell(1, col), lines, fill=WHITE)
for col, lines in enumerate(row_c27):
    set_cell(table.cell(2, col), lines, fill="EAF2FB")

add_text(
    slide,
    "Como ler: Produção OK conta apenas ordens de serviço concluídas com sucesso. Fila de 1,125 dia significa que, "
    "se nenhuma ordem nova entrar, a equipe do Interior precisa de mais de um dia inteiro só para zerar o que já está aberto.",
    56, 416, 1130, 34, size=10.5, color=MUTED,
)

add_rect(slide, 56, 470, 545, 148, NAVY, round_corners=True)
add_text(slide, "DECISÃO OPERACIONAL", 80, 484, 320, 18, size=11, color="9FC9EE", bold=True)
add_text(
    slide,
    "Concentrar a semana em Fortaleza e proteger o Prazo no Interior. A execução começa nominal: "
    "3 áreas e 15 técnicos definem metade da recuperação.",
    80, 506, 498, 100, size=15.5, color=WHITE, bold=True,
)

kpis = [
    ("86 de 127", "técnicos ficaram abaixo da própria meta de produção em 01–03/07"),
    ("245,5 OS", "deixaram de ser executadas — o equivalente a 20 técnicos parados por 3 dias"),
    ("49%", "dessa perda está em apenas 3 das 12 áreas de gestão (GAs)"),
]
for index, (value, label) in enumerate(kpis):
    left = 625 + index * 202
    add_rect(slide, left, 470, 190, 148, WHITE, line=LINE, round_corners=True)
    add_text(slide, value, left + 15, 486, 162, 32, size=23, color=BLUE, bold=True)
    add_text(slide, label, left + 15, 522, 162, 84, size=11, color=MUTED)

# ---------------------------------------------------------------- Slide 2
slide = prs.slides.add_slide(blank)
chrome(slide, 2, 4, "Fonte: Analítico Nominal 01–03/07/2026 e painel de improdutividade de 06/07 · Produtividade OK é a métrica oficial; Baremo OK é visão complementar")
title_block(
    slide,
    "DIAGNÓSTICO DE PRODUÇÃO",
    "O desvio está no time veterano — e boa parte nasce antes de o técnico sair para a rua",
    title_size=30,
)

# painel esquerdo
add_rect(slide, 56, 150, 610, 372, WHITE, line=LINE, round_corners=True)
add_text(slide, "127 técnicos elegíveis", 78, 166, 330, 20, size=15, color=NAVY, bold=True)
add_text(slide, "86 abaixo da meta (68%)", 380, 168, 266, 18, size=13, color=RED, bold=True, align=PP_ALIGN.RIGHT)
add_text(
    slide,
    "Meta individual por tempo de casa: 3 OS/dia até 59 dias, 3,5 de 60 a 89 dias e 4 acima de 90 dias.",
    78, 190, 566, 16, size=10, color=MUTED,
)
segments = [(41, GREEN_DARK), (10, AMBER), (68, ORANGE), (8, RED)]
bar_left, bar_width = 78, 566
offset = bar_left
for count, color in segments:
    seg_width = bar_width * count / 127
    add_rect(slide, offset, 214, seg_width, 20, color)
    offset += seg_width
legend = [
    ("41", "na meta", GREEN_DARK),
    ("10", "novatos abaixo, <90d", AMBER),
    ("68", "veteranos abaixo, >90d", ORANGE),
    ("8", "sem produção", RED),
]
for index, (count, label, color) in enumerate(legend):
    left = 78 + index * 144
    add_rect(slide, left, 244, 7, 28, color)
    add_text(slide, [(count, {"size": 14, "bold": True}), (label, {"size": 8.5, "color": MUTED})],
             left + 12, 242, 132, 34)
add_rect(slide, 78, 284, 566, 42, "FDF4EA")
add_rect(slide, 78, 284, 5, 42, ORANGE)
add_text(
    slide,
    "91% do quadro (115 de 127) é veterano. O problema não é rampa de novato: está no time que já deveria entregar a meta cheia.",
    92, 290, 544, 32, size=10.5, color="5A4632",
)
add_text(slide, "OS não executadas por cluster (gap)", 78, 340, 380, 18, size=13, color=NAVY, bold=True)
add_text(slide, "245,5 OS em 3 dias", 430, 341, 216, 16, size=12, color=RED, bold=True, align=PP_ALIGN.RIGHT)
hbar(slide, "C.27 Fortaleza", "162", 78, 368, 120, 380, 1.0)
hbar(slide, "C.26 Interior", "83,5", 78, 392, 120, 380, 0.515, color="78A9D8")
add_text(
    slide,
    "Gap = soma do que faltou para cada técnico bater a própria meta. Exemplo: meta de 4 OS/dia × 3 dias = 12; "
    "quem entregou 8 deixou 4 OS de gap.",
    78, 420, 566, 30, size=9.5, color=MUTED,
)
add_text(
    slide,
    "As outras 111 visitas do painel ao lado ainda estão sem motivo registrado — reconciliação com o BI até 08/07.",
    78, 456, 566, 30, size=9.5, color=MUTED,
)

# painel direito
add_rect(slide, 686, 150, 538, 372, WHITE, line=LINE, round_corners=True)
add_text(slide, "328 visitas improdutivas", 708, 166, 320, 20, size=15, color=NAVY, bold=True)
add_text(slide, "217 têm motivo e submotivo", 926, 168, 278, 18, size=12, color=RED, bold=True, align=PP_ALIGN.RIGHT)
add_text(
    slide,
    "Visita improdutiva: o técnico foi ao endereço e voltou sem executar o serviço. Motivos registrados nos 4 maiores grupos:",
    708, 190, 494, 30, size=10, color=MUTED,
)
causes = [
    ("OS aberta sem necessidade", "59", 1.0),
    ("Cliente ausente", "52", 0.881),
    ("Cliente pediu reagendamento", "36", 0.61),
    ("Falha massiva na rede", "23", 0.39),
]
for index, (label, value, ratio) in enumerate(causes):
    hbar(slide, label, value, 708, 236 + index * 34, 190, 240, ratio, color=ORANGE)
add_rect(slide, 708, 380, 494, 92, "F1F9F1")
add_rect(slide, 708, 380, 5, 92, GREEN)
add_text(slide, "78%", 724, 392, 90, 30, size=19, color=GREEN_DARK, bold=True)
add_text(
    slide,
    "das improdutivas classificadas nascem antes do despacho: OS mal aberta, cliente não confirmado ou falha de "
    "rede já conhecida. Cobrar só o técnico não resolve.",
    822, 388, 368, 44, size=10.5, color="4A5B52",
)
add_text(
    slide,
    "As outras 111 (34%) não têm motivo no BI — não é mais recorte de tela, é lacuna real do sistema, ainda sem explicação.",
    724, 434, 466, 32, size=9, color="4A5B52",
)

solutions = [
    ("01", "Controle diário em 3 horários",
     "9h: quem está em campo e com carga. 12h: quem ainda não concluiu a primeira OS. 16h30: quem vai fechar o dia "
     "sem produção. Cada caso com motivo e providência no mesmo dia."),
    ("02", "Plano individual por técnico",
     "Prioridade para os 15 casos que definem metade da recuperação: 7 parados, 2 fora da agenda B2C e os maiores "
     "desvios. Detalhe nominal na página seguinte."),
    ("03", "Qualidade da demanda",
     "Travar abertura de OS sem validação técnica, confirmar o cliente antes do despacho e suspender despachos em "
     "áreas com falha massiva ativa."),
]
for index, (number, head, body) in enumerate(solutions):
    left = 56 + index * 398
    add_rect(slide, left, 540, 382, 118, WASH)
    add_rect(slide, left, 540, 382, 3, BLUE)
    add_text(slide, number, left + 14, 552, 40, 26, size=20, color="A2B2C6", bold=True)
    add_text(slide, head, left + 54, 552, 316, 20, size=13, color=NAVY, bold=True)
    add_text(slide, body, left + 54, 574, 316, 78, size=10, color=MUTED)

# ---------------------------------------------------------------- Slide 3
slide = prs.slides.add_slide(blank)
chrome(slide, 3, 4, "Fonte: Analítico Nominal por técnico, 01 a 03/07/2026 · 127 técnicos elegíveis · valores em OS OK")
title_block(
    slide,
    "GESTÃO NOMINAL · 01 A 03/07/2026",
    "A recuperação tem nome: 3 áreas e 15 técnicos concentram metade da perda",
    title_size=30,
)

add_rect(slide, 56, 148, 452, 400, WHITE, line=LINE, round_corners=True)
add_text(slide, "Perda por área de gestão (GA)", 78, 164, 300, 20, size=14, color=NAVY, bold=True)
add_text(slide, "em OS não executadas", 320, 166, 168, 18, size=10, color=MUTED, align=PP_ALIGN.RIGHT)
gas = [
    ("Pedro Henrique · C.27", "63", 1.0, BLUE),
    ("Romildo Albuquerque · C.26", "37", 0.587, BLUE),
    ("Agenor Vieira · C.27", "21", 0.333, BLUE),
    ("Thiago Rabelo · C.27", "20", 0.317, "78A9D8"),
    ("José Orlando · C.27", "19", 0.302, "78A9D8"),
]
for index, (label, value, ratio, color) in enumerate(gas):
    hbar(slide, label, value, 78, 202 + index * 32, 172, 180, ratio, color=color)
add_rect(slide, 78, 372, 408, 152, "F1F9F1")
add_rect(slide, 78, 372, 5, 152, GREEN)
add_text(slide, "26%", 94, 386, 80, 30, size=22, color=GREEN_DARK, bold=True)
add_text(
    slide,
    "de toda a perda do Ceará está numa única área: a de Pedro Henrique, com 21 dos 22 técnicos abaixo da meta. "
    "A mesma área lidera as visitas improdutivas (82 de 328). É por ela que a recuperação começa.",
    94, 418, 376, 100, size=11.5, color="4A5B52",
)

table = add_table(slide, 6, 3, 528, 148, 696, 396, [232, 240, 224])
set_cell(table.cell(0, 0), [("TÉCNICOS", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
set_cell(table.cell(0, 1), [("O QUE OS DADOS MOSTRAM", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
set_cell(table.cell(0, 2), [("AÇÃO IMEDIATA", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
nominal_rows = [
    (
        [("7 técnicos parados", {"size": 10, "bold": True, "color": NAVY}),
         ("Anderson Saraiva, Genilson Nogueira e Felipe Morais (C.26); Emanoel Dias, Egtonio Nunes, Felipe Damasceno e Otávio Dantas (C.27)", {"size": 8, "color": MUTED})],
        [("Nenhum dia trabalhado entre 01 e 03/07. Não é baixa produtividade: é capacidade fora de campo.", {"size": 8.5})],
        [("Confirmar com RH e escala, em 24h, se é férias, atestado ou vaga a repor — e devolver a capacidade ao campo.", {"size": 8.5})],
        WHITE,
    ),
    (
        [("2 técnicos fora da agenda B2C", {"size": 10, "bold": True, "color": NAVY}),
         ("John Wayne Queiroz (C.27) e Jean Carlos Freitas (C.26)", {"size": 8, "color": MUTED})],
        [("Trabalharam os 3 dias, mas quase só em OEM, B2B e apoio — atividades que não contam para a meta. Aparecem como os piores do ranking sem serem.", {"size": 8.5})],
        [("Decidir a alocação: devolver à agenda B2C ou formalizar a função atual e tirá-los desta medição.", {"size": 8.5})],
        "F8FAFD",
    ),
    (
        [("6 técnicos com maior desvio em rota", {"size": 10, "bold": True, "color": NAVY}),
         ("José Danrley (0,7 OS/dia), Jackson Lima (1,3), Daniel Anderson (2,0) — C.26; Matheus Silva (0,5), Antonio Valgleison (2,0), Douglas Maciel (1,7) — C.27", {"size": 8, "color": MUTED})],
        [("Menor entrega entre quem ficou em rota o período todo. Junto com os 2 casos de agenda, os 8 maiores desvios somam 27% da perda.", {"size": 8.5})],
        [("Conversa individual gestor-técnico esta semana, com plano simples de recuperação e acompanhamento diário.", {"size": 8.5})],
        WHITE,
    ),
    (
        [("4 técnicos travados pela demanda", {"size": 10, "bold": True, "color": NAVY}),
         ("André Jefferson, Jadson Araújo, Adriano Sousa e Luiz Nauá (C.27)", {"size": 8, "color": MUTED})],
        [("6 a 7 visitas improdutivas cada um em 3 dias — o dia foi consumido por cliente ausente e OS aberta sem necessidade.", {"size": 8.5})],
        [("Tratar pela frente de qualidade da demanda (confirmação de cliente e trava de abertura); não é caso de cobrança individual.", {"size": 8.5})],
        "F8FAFD",
    ),
    (
        [("Referências a replicar", {"size": 10, "bold": True, "color": GREEN_DARK}),
         ("Carlos André (7,3 OS/dia), Miguel Gomes (7,0) — C.26; Francisco Wagner (6,0) — C.27", {"size": 8, "color": MUTED})],
        [("Quase o dobro da meta, nas mesmas condições de campo. A meta de 4 OS/dia é atingível.", {"size": 8.5})],
        [("Mapear rota, carga e método desses técnicos e usar como padrão nos planos individuais.", {"size": 8.5})],
        "F1F9F1",
    ),
]
for row_index, (col_a, col_b, col_c, fill) in enumerate(nominal_rows, start=1):
    set_cell(table.cell(row_index, 0), col_a, fill=fill, v_margin=4)
    set_cell(table.cell(row_index, 1), col_b, fill=fill, v_margin=4)
    set_cell(table.cell(row_index, 2), col_c, fill=fill, v_margin=4)
for row in table.rows:
    row.height = px(66)

tiers = [
    ("7", "parados — resolver em 24h", RED),
    ("8", "desvio crítico, 27% da perda — plano individual já", RED),
    ("33", "desvio médio (2,5 a 5 OS) — conversa até 10/07", ORANGE),
    ("38", "quase na meta (até 2 OS) — ajuste de carga e rota", AMBER),
    ("41", "na meta — reconhecer e usar de referência", GREEN_DARK),
]
for index, (count, label, color) in enumerate(tiers):
    left = 56 + index * 238
    add_rect(slide, left, 588, 226, 78, WASH)
    add_rect(slide, left, 588, 226, 3, color)
    add_text(slide, count, left + 14, 602, 52, 30, size=22, color=color, bold=True)
    add_text(slide, label, left + 70, 598, 148, 60, size=9.5, color=MUTED)

# ---------------------------------------------------------------- Slide 4
slide = prs.slides.add_slide(blank)
chrome(slide, 4, 4, "Causas não comprovadas pelos dados permanecem marcadas como “em validação”")
title_block(
    slide,
    "FCA · FATO, CAUSA E AÇÃO",
    "Quatro frentes colocam a recuperação em execução nesta semana",
    title_size=30,
)

table = add_table(slide, 5, 6, 42, 148, 1196, 452, [200, 258, 268, 116, 216, 138])
fca_heads = ["FATO", "CAUSA", "AÇÃO", "PRAZO", "RESPONSÁVEL", "STATUS"]
for col, head in enumerate(fca_heads):
    set_cell(table.cell(0, col), [(head, {"size": 10.5, "bold": True, "color": WHITE})], fill=NAVY)

fca_rows = [
    (
        [("Fortaleza abaixo nos 3 Ps", {"size": 11, "bold": True, "color": NAVY}),
         ("Presença 94,5% · Produção 3,2 OS/dia · Prazo 83,2%", {"size": 9, "color": MUTED})],
        [("Em Produção, 61 dos 81 técnicos estão abaixo da meta, com metade da perda em 3 áreas (Pedro Henrique, Agenor e Thiago Rabelo). Para Presença e Prazo a causa ainda está em levantamento.", {"size": 9.5})],
        [("Controle diário em 3 horários — 9h presença e carga, 12h primeira OS concluída, 16h30 risco de dia sem produção — com lista nominal e providência no mesmo dia.", {"size": 9.5})],
        [("07/07", {"size": 11, "bold": True, "color": NAVY}), ("1ª revisão 13/07", {"size": 9, "color": MUTED})],
        [("José Nilton (GO), Jefferson Oliveira e GAs de Fortaleza", {"size": 9.5})],
        ("A iniciar", BLUE_LIGHT, "114F88"),
        WHITE,
    ),
    (
        [("245,5 OS não executadas", {"size": 11, "bold": True, "color": NAVY}),
         ("86 dos 127 técnicos abaixo da meta em 01–03/07", {"size": 9, "color": MUTED})],
        [("Três perfis distintos: 7 técnicos sem nenhum dia trabalhado, 2 alocados fora da agenda B2C e 8 casos críticos que somam 27% da perda; os demais estão a menos de 5 OS da meta.", {"size": 9.5})],
        [("Confirmar com RH a situação dos 7 parados; decidir a alocação dos 2 desviados; plano individual para os 8 críticos; ajuste de carga e rota para os demais.", {"size": 9.5})],
        [("08/07", {"size": 11, "bold": True, "color": NAVY}), ("planos até 10/07", {"size": 9, "color": MUTED})],
        [("GAs dos dois clusters, coordenados pelos GOs", {"size": 9.5})],
        ("Imediato", RED_LIGHT, "8A1717"),
        "F8FAFD",
    ),
    (
        [("328 visitas improdutivas", {"size": 11, "bold": True, "color": NAVY}),
         ("217 têm motivo e submotivo registrado", {"size": 9, "color": MUTED})],
        [("Entre as 217 classificadas, 78% nascem antes do despacho: OS aberta sem necessidade, cliente não confirmado, reagendamento não retirado da fila e falha massiva. As outras 111 (34%) não têm motivo no BI — lacuna real, não recorte de tela.", {"size": 9.5})],
        [("Reconciliar com o BI por que 111 baixas não têm motivo registrado; travar abertura de OS sem validação técnica; confirmar o cliente antes do despacho; suspender despacho em área com falha massiva ativa.", {"size": 9.5})],
        [("08/07", {"size": 11, "bold": True, "color": NAVY}), ("travas até 10/07", {"size": 9, "color": MUTED})],
        [("Biondillo (Performance), Davi dos Reis Luz (BI) e GOs do Ceará", {"size": 9.5})],
        ("Em validação", "FFF0C7", "765400"),
        WHITE,
    ),
    (
        [("Fila do Interior acima do limite", {"size": 11, "bold": True, "color": NAVY}),
         ("1,125 dia de trabalho acumulado em 06/07", {"size": 9, "color": MUTED})],
        [("Ainda sem dados de entrada de ordens por hora — não é possível afirmar se é excesso de demanda ou perda de capacidade. Causa em validação.", {"size": 9.5})],
        [("Toda manhã, atacar primeiro as ordens abertas há mais de 20 horas, antes de estourarem o prazo de 24h, e medir a fila no fim do dia até voltar a 1 dia ou menos.", {"size": 9.5})],
        [("07/07", {"size": 11, "bold": True, "color": NAVY}), ("fila ≤ 1 dia até 10/07", {"size": 9, "color": MUTED})],
        [("José Geraldo (GO), Francisco Cleiton e GAs do Interior", {"size": 9.5})],
        ("Atenção", "FFE5CB", "8A4C00"),
        "F8FAFD",
    ),
]
for row_index, (fato, causa, acao, prazo, resp, status, fill) in enumerate(fca_rows, start=1):
    set_cell(table.cell(row_index, 0), fato, fill=fill)
    set_cell(table.cell(row_index, 1), causa, fill=fill)
    set_cell(table.cell(row_index, 2), acao, fill=fill)
    set_cell(table.cell(row_index, 3), prazo, fill=fill)
    set_cell(table.cell(row_index, 4), resp, fill=fill)
    status_text, status_fill, status_color = status
    set_cell(
        table.cell(row_index, 5),
        [(status_text, {"size": 10, "bold": True, "color": status_color, "align": PP_ALIGN.CENTER})],
        fill=status_fill,
    )

add_rect(slide, 42, 618, 1196, 48, NAVY)
add_rect(slide, 42, 618, 7, 48, GREEN)
add_text(slide, "DECISÕES SOLICITADAS", 64, 634, 200, 18, size=11.5, color=WHITE, bold=True)
add_text(
    slide,
    "Confirmar responsáveis e prazos   ·   Ratificar 85% como meta do Prazo 24h   ·   Manter Produtividade OK como métrica oficial",
    270, 634, 940, 18, size=11.5, color="D5DEF3",
)

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(f"OK: {OUTPUT}")
