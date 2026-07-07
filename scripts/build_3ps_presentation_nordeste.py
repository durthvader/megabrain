"""Gera outputs/3ps/3Ps_Nordeste_Executivo.pptx com a identidade visual Alloha.

Espelha build_3ps_presentation_mapi.py, adaptado para a regional R7.1 (Nordeste: AL, BA, PB, PE, RN, SE).
Uso: python scripts/build_3ps_presentation_nordeste.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "outputs" / "3ps" / "3Ps_Nordeste_Executivo.pptx"

NAVY = "081551"
BLUE = "1476C9"
BLUE_LIGHT = "DCECF9"
GREEN = "42DF4B"
GREEN_DARK = "13853A"
RED = "D94545"
RED_LIGHT = "FDE2DF"
ORANGE = "F08A24"
AMBER = "E4BD30"
INK = "182238"
MUTED = "667085"
LINE = "D9E0EA"
WASH = "F3F6FA"
WHITE = "FFFFFF"


def px(value: float) -> Emu:
    return Emu(int(value * 9525))


def rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


def add_rect(slide, left, top, width, height, fill, line=None, line_width=1.0, round_corners=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_corners else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, px(left), px(top), px(width), px(height))
    if round_corners:
        shape.adjustments[0] = 0.07
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_text(slide, value, left, top, width, height, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Aptos", line_spacing=1.06):
    box = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, margin, 0)
    lines = value if isinstance(value, list) else [(value, {})]
    for index, (text, opts) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = opts.get("align", align)
        paragraph.line_spacing = opts.get("line_spacing", line_spacing)
        if index > 0:
            paragraph.space_before = Pt(opts.get("space_before", 2))
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        run.font.color.rgb = rgb(opts.get("color", color))
        run.font.name = opts.get("font", font)
    return box


def chrome(slide, number, total, source):
    add_rect(slide, 0, 0, 1100, 8, NAVY)
    add_rect(slide, 1100, 0, 180, 8, GREEN)
    add_text(slide, "alloha", 1090, 30, 130, 36, size=27, color=NAVY,
             align=PP_ALIGN.RIGHT, font="Aptos Display")
    add_text(slide, "F I B R A", 1090, 66, 130, 14, size=8, color=GREEN_DARK,
             bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, source, 56, 688, 1040, 16, size=9, color="8B94A6")
    add_text(slide, f"{number} / {total}", 1150, 688, 74, 16, size=9,
             color="8B94A6", align=PP_ALIGN.RIGHT)


def title_block(slide, eyebrow, title, lead=None, title_size=34, lead_top=138):
    add_text(slide, eyebrow, 56, 36, 900, 20, size=12, color=BLUE, bold=True)
    add_text(slide, title, 56, 60, 1000, 76, size=title_size, color=NAVY,
             bold=True, font="Aptos Display", line_spacing=1.0)
    if lead:
        add_text(slide, lead, 56, lead_top, 1120, 44, size=15, color="45516A")


def set_cell(cell, lines, fill=None, anchor=MSO_ANCHOR.MIDDLE, v_margin=6):
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill)
    cell.vertical_anchor = anchor
    cell.margin_left = px(10)
    cell.margin_right = px(8)
    cell.margin_top = px(v_margin)
    cell.margin_bottom = px(v_margin)
    tf = cell.text_frame
    tf.word_wrap = True
    for index, (text, opts) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = opts.get("align", PP_ALIGN.LEFT)
        paragraph.line_spacing = opts.get("line_spacing", 1.05)
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", 11))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = rgb(opts.get("color", INK))
        run.font.name = opts.get("font", "Aptos")


def add_table(slide, rows, cols, left, top, width, height, col_widths):
    frame = slide.shapes.add_table(rows, cols, px(left), px(top), px(width), px(height))
    table = frame.table
    table.first_row = False
    table.horz_banding = False
    for index, cw in enumerate(col_widths):
        table.columns[index].width = px(cw)
    return table


def hbar(slide, label, value_text, left, top, label_width, bar_width, ratio,
         color=BLUE, track="E7EDF5", bar_height=12):
    add_text(slide, label, left, top - 2, label_width, 16, size=10.5, color=INK)
    track_left = left + label_width + 8
    add_rect(slide, track_left, top, bar_width, bar_height, track)
    add_rect(slide, track_left, top, max(bar_width * ratio, 4), bar_height, color)
    add_text(slide, value_text, track_left + bar_width + 8, top - 2, 44, 16,
             size=11, color=INK, bold=True)


prs = Presentation()
prs.slide_width = px(1280)
prs.slide_height = px(720)
blank = prs.slide_layouts[6]

# ---------------------------------------------------------------- Slide 1
slide = prs.slides.add_slide(blank)
chrome(slide, 1, 4, "Fonte: exportações de BI operacional da regional R7.1 (AL/BA/PB/PE/RN/SE), 01–07/07/2026 · OS = ordem de serviço concluída com baixa OK")
title_block(
    slide,
    "OPERAÇÕES · NORDESTE (R7.1) · 07/07/2026",
    "Sergipe e Pernambuco fecham 06/07 acima da meta de Produção e Prazo, mas a produtividade da semana caiu 32% "
    "e Sergipe está abaixo na Agenda",
    "No fechamento de 06/07, Sergipe (C.21) e Pernambuco (C.23) superaram a meta de Produção e de Prazo 24h; "
    "Pernambuco também está acima da meta de Cumprimento de Agenda, enquanto Sergipe fica 5 p.p. abaixo. Na "
    "visão semanal, a Produção 11 OK da regional caiu de 4,4 para 3,0 OS/dia (-32%) entre a semana fechada (S27) "
    "e a semana em curso (S28).",
    title_size=27,
    lead_top=152,
)

table = add_table(slide, 3, 4, 56, 202, 1168, 204, [230, 306, 306, 326])
headers = [
    ("Cluster", ""),
    ("Baixando serviço", "vs. folha meta, 06/07"),
    ("Produção", "% da meta esperada, 06/07"),
    ("Prazo 24h", "meta 85% das ordens"),
]
for col, (head, sub) in enumerate(headers):
    lines = [(head, {"size": 13, "bold": True, "color": WHITE})]
    if sub:
        lines.append((sub, {"size": 9.5, "color": "BFC9E3"}))
    set_cell(table.cell(0, col), lines, fill=NAVY)

row_c21 = [
    [("C.21", {"size": 16, "bold": True, "color": NAVY}), ("Sergipe", {"size": 10, "color": MUTED})],
    [("87,5%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("14 de 16 da folha meta", {"size": 10, "color": MUTED})],
    [("118,1%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("acima do esperado", {"size": 10, "color": MUTED})],
    [("90,0%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("5 p.p. acima da meta", {"size": 10, "color": MUTED})],
]
row_c23 = [
    [("C.23", {"size": 16, "bold": True, "color": NAVY}), ("Pernambuco", {"size": 10, "color": MUTED})],
    [("94,1%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("16 de 17 da folha meta", {"size": 10, "color": MUTED})],
    [("110,7%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("acima do esperado", {"size": 10, "color": MUTED})],
    [("87,5%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("2,5 p.p. acima da meta", {"size": 10, "color": MUTED})],
]
for col, lines in enumerate(row_c21):
    set_cell(table.cell(1, col), lines, fill=WHITE)
for col, lines in enumerate(row_c23):
    set_cell(table.cell(2, col), lines, fill="EAF2FB")

add_text(
    slide,
    "Como ler: Baixando Serviço mede quantos técnicos da folha meta estão de fato executando ordens no dia. A "
    "regional R7.1 fechou 06/07 com Prazo 24h em 88,9% e Cumprimento de Agenda em 86,1% — ambos acima da meta de "
    "85%, puxados por Pernambuco. Sergipe fica 5 p.p. abaixo da meta de Agenda (80,0%); Pernambuco está 8,8 p.p. "
    "acima (93,8%).",
    56, 416, 1130, 34, size=10.5, color=MUTED,
)

add_rect(slide, 56, 470, 545, 168, NAVY, round_corners=True)
add_text(slide, "DECISÃO OPERACIONAL", 80, 484, 320, 18, size=11, color="9FC9EE", bold=True)
add_text(
    slide,
    "Reverter a queda de produtividade da S28 antes que vire tendência e fechar a lacuna de Cumprimento de Agenda "
    "em Sergipe (80% vs meta de 85%). Execução também nominal: a GA de Leandro Sousa de Lorena concentra o maior "
    "gap em Pernambuco, e 8 técnicos concentram os maiores desvios.",
    80, 506, 498, 118, size=14, color=WHITE, bold=True,
)

kpis = [
    ("6 de 64", "registros-semana em Q4/NP (produtividade mais baixa ou sem produção) no quartil 11 OK"),
    ("63 OS", "de gap estimado em 7 dias (01–07/07), considerando meta por tempo de casa"),
    ("-32%", "queda de Produção 11 OK entre a semana fechada e a semana em curso"),
]
for index, (value, label) in enumerate(kpis):
    left = 625 + index * 202
    add_rect(slide, left, 470, 190, 168, WHITE, line=LINE, round_corners=True)
    add_text(slide, value, left + 15, 486, 162, 32, size=23, color=BLUE, bold=True)
    add_text(slide, label, left + 15, 522, 162, 104, size=11, color=MUTED)

# ---------------------------------------------------------------- Slide 2
slide = prs.slides.add_slide(blank)
chrome(slide, 2, 4, "Fonte: Analítico Nominal 01–07/07/2026, Quartil de Produtividade S27/S28 e exportações de motivo/submotivo de improdutividade da R7.1")
title_block(
    slide,
    "DIAGNÓSTICO DE PRODUÇÃO",
    "O quadro está concentrado nos quartis mais produtivos — o desafio é sustentar o ritmo da S27 na semana em curso",
    title_size=25,
)

add_rect(slide, 56, 150, 610, 372, WHITE, line=LINE, round_corners=True)
add_text(slide, "32 técnicos elegíveis", 78, 166, 330, 20, size=15, color=NAVY, bold=True)
add_text(slide, "Quadro 100% veterano", 340, 168, 306, 18, size=12, color=MUTED, align=PP_ALIGN.RIGHT)
add_text(
    slide,
    "Quartil de produtividade (11 OK), S27 e S28. Todo o quadro tem mais de 90 dias de casa: 16 técnicos em "
    "Sergipe, 16 em Pernambuco.",
    78, 190, 566, 28, size=10, color=MUTED,
)
segments = [(43.8, GREEN_DARK), (25.0, AMBER), (21.9, ORANGE), (9.4, RED)]
bar_left, bar_width = 78, 566
offset = bar_left
for pct, color in segments:
    seg_width = bar_width * pct / 100
    add_rect(slide, offset, 222, seg_width, 20, color)
    offset += seg_width
legend = [
    ("43,8%", "Q1 — produtividade mais alta", GREEN_DARK),
    ("25,0%", "Q2 — intermediária", AMBER),
    ("21,9%", "Q3 — abaixo da intermediária", ORANGE),
    ("9,4%", "Q4/NP — mais baixa ou sem produção", RED),
]
for index, (count, label, color) in enumerate(legend):
    left = 78 + index * 144
    add_rect(slide, left, 252, 7, 28, color)
    add_text(slide, [(count, {"size": 14, "bold": True}), (label, {"size": 8.5, "color": MUTED})],
             left + 12, 250, 132, 34)
add_rect(slide, 78, 292, 566, 46, "F1F9F1")
add_rect(slide, 78, 292, 5, 46, GREEN)
add_text(
    slide,
    "O quadro está concentrado nos quartis de maior produtividade — quase 70% dos registros-semana estão em "
    "Q1/Q2. O risco não está na distribuição do time: está na queda de ritmo entre a semana fechada e a atual.",
    92, 298, 544, 36, size=10.5, color="2F5C3E",
)
add_text(slide, "Gap estimado por cluster (01–07/07)", 78, 350, 380, 18, size=13, color=NAVY, bold=True)
add_text(slide, "63 OS em 7 dias", 430, 351, 216, 16, size=12, color=RED, bold=True, align=PP_ALIGN.RIGHT)
hbar(slide, "C.23 Pernambuco", "40", 78, 378, 120, 380, 1.0)
hbar(slide, "C.21 Sergipe", "23", 78, 402, 120, 380, 0.575, color="78A9D8")
add_text(
    slide,
    "Gap = máximo(meta por faixa de tempo de casa × dias trabalhados − volume OK, 0), estimado a partir do "
    "Analítico Nominal (dias trabalhados inferidos por Volume OK / Prod. OK).",
    78, 430, 566, 30, size=9.5, color=MUTED,
)
add_text(
    slide,
    "1 técnico (Ramon dos Santos, C.21) ficou sem nenhuma baixa no período inteiro — confirmar com RH/escala.",
    78, 462, 566, 30, size=9.5, color=MUTED,
)

add_rect(slide, 686, 150, 538, 372, WHITE, line=LINE, round_corners=True)
add_text(slide, "98 visitas improdutivas", 708, 166, 320, 20, size=15, color=NAVY, bold=True)
add_text(slide, "82 com motivo registrado", 866, 168, 338, 18, size=12, color=RED, bold=True, align=PP_ALIGN.RIGHT)
add_text(
    slide,
    "Motivos registrados nos 4 maiores grupos (01–07/07):",
    708, 190, 494, 30, size=10, color=MUTED,
)
causes = [
    ("Cliente ausente", "19", 1.0),
    ("Abertura indevida", "15", 0.789),
    ("Solicitação de reagendamento", "12", 0.632),
    ("Falha massiva", "11", 0.579),
]
for index, (label, value, ratio) in enumerate(causes):
    hbar(slide, label, value, 708, 236 + index * 34, 190, 240, ratio, color=ORANGE)
add_rect(slide, 708, 380, 494, 92, "F1F9F1")
add_rect(slide, 708, 380, 5, 92, GREEN)
add_text(slide, "70%", 724, 392, 90, 30, size=19, color=GREEN_DARK, bold=True)
add_text(
    slide,
    "das improdutivas classificadas estão nos 4 principais motivos, com “Cliente ausente” e “Abertura indevida” "
    "à frente — parte relevante nasce antes do despacho ou por ausência do cliente na visita.",
    822, 388, 368, 48, size=10.5, color="4A5B52",
)
solutions = [
    ("01", "Reverter a queda S27→S28",
     "Produção 11 OK caiu de 4,4 para 3,0 OS/dia. Acompanhamento diário para retomar o patamar da semana "
     "anterior até o fim da S28."),
    ("02", "Plano individual nominal",
     "Foco nos 8 maiores gaps individuais e na GA de Pernambuco com maior concentração de perda. Detalhe na "
     "página seguinte."),
    ("03", "Qualidade da demanda",
     "Travar abertura de OS sem validação técnica e confirmar a presença do cliente antes do despacho — os dois "
     "motivos somam 34 das 82 baixas classificadas."),
]
for index, (number, head, body) in enumerate(solutions):
    left = 56 + index * 398
    add_rect(slide, left, 540, 382, 118, WASH)
    add_rect(slide, left, 540, 382, 3, BLUE)
    add_text(slide, number, left + 14, 552, 40, 26, size=20, color="A2B2C6", bold=True)
    add_text(slide, head, left + 54, 552, 316, 20, size=13, color=NAVY, bold=True)
    add_text(slide, body, left + 54, 574, 316, 78, size=10, color=MUTED)

# ---------------------------------------------------------------- Slide 3
slide = prs.slides.add_slide(blank)
chrome(slide, 3, 4, "Fonte: Analítico Nominal por técnico, 01 a 07/07/2026 · 32 técnicos elegíveis · gap estimado, dias trabalhados inferidos")
title_block(
    slide,
    "GESTÃO NOMINAL · 01 A 07/07/2026",
    "Uma área de gestão em Pernambuco e oito técnicos concentram a maior parte do gap da regional",
    title_size=25,
)

add_rect(slide, 56, 148, 452, 400, WHITE, line=LINE, round_corners=True)
add_text(slide, "Gap por área de gestão (GA)", 78, 164, 300, 20, size=14, color=NAVY, bold=True)
add_text(slide, "em OS não executadas (estimado)", 300, 166, 188, 18, size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)
gas = [
    ("Leandro Sousa de Lorena · C.23", "37", 1.0, BLUE),
    ("Caio Vinicius Rosendo dos Santos · C.21", "23", 0.622, BLUE),
    ("Luiz Cavalcanti da Silva Neto · C.23", "3", 0.081, "78A9D8"),
]
for index, (label, value, ratio, color) in enumerate(gas):
    hbar(slide, label, value, 78, 202 + index * 32, 172, 180, ratio, color=color)
add_rect(slide, 78, 366, 408, 168, "F1F9F1")
add_rect(slide, 78, 366, 5, 168, GREEN)
add_text(slide, "1 técnico", 94, 378, 120, 30, size=18, color=GREEN_DARK, bold=True)
add_text(
    slide,
    "ficou sem nenhuma baixa no período inteiro (01–07/07): Ramon dos Santos (C.21, GA Caio Vinicius Rosendo dos "
    "Santos) — capacidade fora de campo a confirmar com RH/escala.",
    94, 410, 376, 116, size=10.5, color="4A5B52",
)

table = add_table(slide, 4, 3, 528, 148, 696, 396, [232, 240, 224])
set_cell(table.cell(0, 0), [("TÉCNICOS", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
set_cell(table.cell(0, 1), [("O QUE OS DADOS MOSTRAM", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
set_cell(table.cell(0, 2), [("AÇÃO IMEDIATA", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
nominal_rows = [
    (
        [("1 técnico sem baixa no período", {"size": 10, "bold": True, "color": NAVY}),
         ("Ramon dos Santos (C.21)", {"size": 8, "color": MUTED})],
        [("Nenhuma OS concluída entre 01 e 07/07. Não é baixa produtividade: é capacidade fora de campo.", {"size": 8.5})],
        [("Confirmar com RH e escala em 24h se é férias, atestado ou vaga a repor.", {"size": 8.5})],
        WHITE,
    ),
    (
        [("8 maiores desvios individuais", {"size": 10, "bold": True, "color": NAVY}),
         ("Idecacio Lucas dos Santos (gap 8); Anderson Henrique da Silva Ferreira, Heitor Cezar Costa da Silva e Tiago de Jesus Andrade (gap 6); Fabio Santos Silva e Andre Felipe Luis (gap 5); Ubiratan Rodrigues de Barros e Darlan de Oliveira Rocha (gap 4) — C.21/C.23", {"size": 8, "color": MUTED})],
        [("Concentram 44 das 63 OS de gap da regional; a maioria com poucos dias trabalhados no período.", {"size": 8.5})],
        [("Conversa individual gestor-técnico esta semana, com plano simples de recuperação e verificação da escala real.", {"size": 8.5})],
        "F8FAFD",
    ),
    (
        [("Referências a replicar", {"size": 10, "bold": True, "color": GREEN_DARK}),
         ("Edilson Pereira da Silva (6,0) e Jose Paulo de Sousa (5,67) — C.23; Gustavo da Silva Bonifacio (5,6), Erlan Cardoso da Silva (4,6) e Emanoel Santos Vieira (4,4) — C.21", {"size": 8, "color": MUTED})],
        [("Bem acima da meta de 4 OS/dia, nas mesmas condições de campo.", {"size": 8.5})],
        [("Mapear rota, carga e método desses técnicos e usar como padrão, priorizando a GA de Leandro Sousa de Lorena.", {"size": 8.5})],
        "F1F9F1",
    ),
]
for row_index, (col_a, col_b, col_c, fill) in enumerate(nominal_rows, start=1):
    set_cell(table.cell(row_index, 0), col_a, fill=fill, v_margin=4)
    set_cell(table.cell(row_index, 1), col_b, fill=fill, v_margin=4)
    set_cell(table.cell(row_index, 2), col_c, fill=fill, v_margin=4)
for row in table.rows:
    row.height = px(88)

tiers = [
    ("1", "sem baixa no período — resolver em 24h", RED),
    ("8", "maior desvio individual — plano já", RED),
    ("28", "Q1 — produtividade mais alta, usar de referência", GREEN_DARK),
    ("16", "Q2 — produtividade intermediária", AMBER),
    ("20", "Q3/Q4/NP — abaixo do ideal ou sem produção", ORANGE),
]
for index, (count, label, color) in enumerate(tiers):
    left = 56 + index * 238
    add_rect(slide, left, 588, 226, 78, WASH)
    add_rect(slide, left, 588, 226, 3, color)
    add_text(slide, count, left + 14, 602, 52, 30, size=22, color=color, bold=True)
    add_text(slide, label, left + 70, 598, 148, 60, size=9.5, color=MUTED)

# ---------------------------------------------------------------- Slide 4
slide = prs.slides.add_slide(blank)
chrome(slide, 4, 4, "Causas não comprovadas pelos dados permanecem marcadas como “em validação”")
title_block(
    slide,
    "FCA · FATO, CAUSA E AÇÃO",
    "Quatro frentes colocam a recuperação do Nordeste em execução nesta semana",
    title_size=27,
)

table = add_table(slide, 5, 6, 42, 148, 1196, 452, [200, 258, 268, 116, 216, 138])
fca_heads = ["FATO", "CAUSA", "AÇÃO", "PRAZO", "RESPONSÁVEL", "STATUS"]
for col, head in enumerate(fca_heads):
    set_cell(table.cell(0, col), [(head, {"size": 10.5, "bold": True, "color": WHITE})], fill=NAVY)

fca_rows = [
    (
        [("Sergipe abaixo na meta de Cumprimento de Agenda", {"size": 11, "bold": True, "color": NAVY}),
         ("Agenda 80,0% em Sergipe · 93,8% em Pernambuco, em 06/07", {"size": 9, "color": MUTED})],
        [("Produção e Prazo 24h seguem dentro da meta nos dois clusters — o desvio está concentrado na agenda de Sergipe.", {"size": 9.5})],
        [("Controle diário de Cumprimento de Agenda em Sergipe, com lista nominal de OS não atendidas e providência no mesmo dia.", {"size": 9.5})],
        [("08/07", {"size": 11, "bold": True, "color": NAVY}), ("1ª revisão 15/07", {"size": 9, "color": MUTED})],
        [("Gestores de operação de Sergipe", {"size": 9.5})],
        ("A iniciar", BLUE_LIGHT, "114F88"),
        WHITE,
    ),
    (
        [("63 OS de gap estimado em 7 dias", {"size": 11, "bold": True, "color": NAVY}),
         ("concentrado na GA de Leandro Sousa de Lorena (PE) e 8 técnicos críticos", {"size": 9, "color": MUTED})],
        [("1 técnico sem nenhuma baixa no período; 8 com maior desvio individual, a maioria com poucos dias trabalhados.", {"size": 9.5})],
        [("Confirmar com RH a situação do técnico parado; plano individual para os 8 críticos; ajuste de carga e rota para os demais.", {"size": 9.5})],
        [("08/07", {"size": 11, "bold": True, "color": NAVY}), ("planos até 10/07", {"size": 9, "color": MUTED})],
        [("GA Leandro Sousa de Lorena e GA Caio Vinicius Rosendo dos Santos", {"size": 9.5})],
        ("Imediato", RED_LIGHT, "8A1717"),
        "F8FAFD",
    ),
    (
        [("98 visitas improdutivas no período", {"size": 11, "bold": True, "color": NAVY}),
         ("82 com motivo registrado", {"size": 9, "color": MUTED})],
        [("Cliente ausente (19) e Abertura indevida (15) lideram os motivos registrados — 70% das classificadas concentradas nos 4 principais motivos.", {"size": 9.5})],
        [("Confirmar o cliente antes do despacho; travar abertura de OS sem validação técnica; revisar agendamento para reduzir reagendamento.", {"size": 9.5})],
        [("08/07", {"size": 11, "bold": True, "color": NAVY}), ("travas até 10/07", {"size": 9, "color": MUTED})],
        [("Equipe de BI/Performance regional", {"size": 9.5})],
        ("Em validação", "FFF0C7", "765400"),
        WHITE,
    ),
    (
        [("Queda de 32% na Produção 11 OK (S27→S28)", {"size": 11, "bold": True, "color": NAVY}),
         ("4,4 → 3,0 OS/dia", {"size": 9, "color": MUTED})],
        [("Queda presente nos dois clusters, mais acentuada em Sergipe (-40%) do que em Pernambuco (-26%).", {"size": 9.5})],
        [("Acompanhar o fechamento da S28 frente à S27, priorizando as GAs com maior gap identificado (Leandro Sousa de Lorena e Caio Vinicius Rosendo dos Santos).", {"size": 9.5})],
        [("12/07", {"size": 11, "bold": True, "color": NAVY}), ("fechamento da S28", {"size": 9, "color": MUTED})],
        [("Gestores de operação da regional R7.1", {"size": 9.5})],
        ("Atenção", "FFE5CB", "8A4C00"),
        "F8FAFD",
    ),
]
for row_index, (fato, causa, acao, prazo, resp, status, fill) in enumerate(fca_rows, start=1):
    set_cell(table.cell(row_index, 0), fato, fill=fill)
    set_cell(table.cell(row_index, 1), causa, fill=fill)
    set_cell(table.cell(row_index, 2), acao, fill=fill)
    set_cell(table.cell(row_index, 3), prazo, fill=fill)
    set_cell(table.cell(row_index, 4), resp, fill=fill)
    status_text, status_fill, status_color = status
    set_cell(
        table.cell(row_index, 5),
        [(status_text, {"size": 10, "bold": True, "color": status_color, "align": PP_ALIGN.CENTER})],
        fill=status_fill,
    )

add_rect(slide, 42, 618, 1196, 48, NAVY)
add_rect(slide, 42, 618, 7, 48, GREEN)
add_text(slide, "DECISÕES SOLICITADAS", 64, 634, 200, 18, size=11.5, color=WHITE, bold=True)
add_text(
    slide,
    "Confirmar responsáveis e prazos   ·   Ratificar 85% como meta de Prazo 24h e Cumprimento de Agenda   ·   "
    "Priorizar a GA de Leandro Sousa de Lorena na recuperação",
    270, 634, 940, 18, size=10.5, color="D5DEF3",
)

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(f"OK: {OUTPUT}")
