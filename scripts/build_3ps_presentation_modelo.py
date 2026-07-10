"""Gera outputs/3ps/3Ps_Modelo_Executivo.pptx — MODELO editável para as regionais.

Mesmo layout e identidade Alloha das apresentações do Ceará, Nordeste e MA/PI,
porém com conteúdo de EXEMPLO neutro. A regional abre o .pptx no PowerPoint e
sobrescreve os textos direto nas caixas e tabelas. Cada slide traz instruções de
preenchimento nas NOTAS (Exibir > Anotações), que não aparecem na apresentação.

Uso: python scripts/build_3ps_presentation_modelo.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "outputs" / "3ps" / "3Ps_Modelo_Executivo.pptx"

NAVY = "081551"
BLUE = "1476C9"
BLUE_LIGHT = "DCECF9"
GREEN = "42DF4B"
GREEN_DARK = "13853A"   # texto/estado "na meta" (OK)
RED = "D94545"          # texto/estado "abaixo da meta" (RISCO)
RED_LIGHT = "FDE2DF"
ORANGE = "F08A24"
AMBER = "E4BD30"
WARN = "9A6200"         # texto/estado "atenção / no limite"
INK = "182238"
MUTED = "667085"
LINE = "D9E0EA"
WASH = "F3F6FA"
WHITE = "FFFFFF"


def px(value: float) -> Emu:
    return Emu(int(value * 9525))


def rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


def add_rect(slide, left, top, width, height, fill, line=None, line_width=1.0, round_corners=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_corners else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, px(left), px(top), px(width), px(height))
    if round_corners:
        shape.adjustments[0] = 0.07
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_text(slide, value, left, top, width, height, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Aptos", line_spacing=1.06):
    """value: str ou lista de (texto, opções por linha)."""
    box = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, margin, 0)
    lines = value if isinstance(value, list) else [(value, {})]
    for index, (text, opts) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = opts.get("align", align)
        paragraph.line_spacing = opts.get("line_spacing", line_spacing)
        if index > 0:
            paragraph.space_before = Pt(opts.get("space_before", 2))
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        run.font.color.rgb = rgb(opts.get("color", color))
        run.font.name = opts.get("font", font)
    return box


def add_notes(slide, text):
    """Instruções de preenchimento — visíveis só no modo de edição (Anotações)."""
    slide.notes_slide.notes_text_frame.text = text


def chrome(slide, number, total, source):
    add_rect(slide, 0, 0, 1100, 8, NAVY)
    add_rect(slide, 1100, 0, 180, 8, GREEN)
    add_text(slide, "alloha", 1090, 30, 130, 36, size=27, color=NAVY,
             align=PP_ALIGN.RIGHT, font="Aptos Display")
    add_text(slide, "F I B R A", 1090, 66, 130, 14, size=8, color=GREEN_DARK,
             bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, source, 56, 688, 1040, 16, size=9, color="8B94A6")
    add_text(slide, f"{number} / {total}", 1150, 688, 74, 16, size=9,
             color="8B94A6", align=PP_ALIGN.RIGHT)


def title_block(slide, eyebrow, title, lead=None, title_size=32, lead_top=138):
    add_text(slide, eyebrow, 56, 36, 900, 20, size=12, color=BLUE, bold=True)
    add_text(slide, title, 56, 60, 1000, 76, size=title_size, color=NAVY,
             bold=True, font="Aptos Display", line_spacing=1.0)
    if lead:
        add_text(slide, lead, 56, lead_top, 1120, 44, size=15, color="45516A")


def set_cell(cell, lines, fill=None, anchor=MSO_ANCHOR.MIDDLE, v_margin=6):
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill)
    cell.vertical_anchor = anchor
    cell.margin_left = px(10)
    cell.margin_right = px(8)
    cell.margin_top = px(v_margin)
    cell.margin_bottom = px(v_margin)
    tf = cell.text_frame
    tf.word_wrap = True
    for index, (text, opts) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = opts.get("align", PP_ALIGN.LEFT)
        paragraph.line_spacing = opts.get("line_spacing", 1.05)
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", 11))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = rgb(opts.get("color", INK))
        run.font.name = opts.get("font", "Aptos")


def add_table(slide, rows, cols, left, top, width, height, col_widths):
    frame = slide.shapes.add_table(rows, cols, px(left), px(top), px(width), px(height))
    table = frame.table
    table.first_row = False
    table.horz_banding = False
    for index, cw in enumerate(col_widths):
        table.columns[index].width = px(cw)
    return table


def hbar(slide, label, value_text, left, top, label_width, bar_width, ratio,
         color=BLUE, track="E7EDF5", bar_height=12):
    add_text(slide, label, left, top - 2, label_width, 16, size=10.5, color=INK)
    track_left = left + label_width + 8
    add_rect(slide, track_left, top, bar_width, bar_height, track)
    add_rect(slide, track_left, top, max(bar_width * ratio, 4), bar_height, color)
    add_text(slide, value_text, track_left + bar_width + 8, top - 2, 44, 16,
             size=11, color=INK, bold=True)


prs = Presentation()
prs.slide_width = px(1280)
prs.slide_height = px(720)
blank = prs.slide_layouts[6]

# ================================================================ Slide 1
slide = prs.slides.add_slide(blank)
chrome(slide, 1, 4, "Fonte: [descreva a base e o período] · OS = ordem de serviço concluída com baixa OK")
title_block(
    slide,
    "OPERAÇÕES · [REGIÃO] ([CÓDIGO]) · [DD/MM/AAAA]",
    "[Manchete da regional: a mensagem principal do fechamento em uma frase]",
    "[Resumo em 2 a 3 linhas: como cada cluster fechou nos indicadores, destacando o que está "
    "na meta e o que ficou abaixo. Este é o texto de contexto que a diretoria lê primeiro.]",
    title_size=30,
    lead_top=152,
)

table = add_table(slide, 3, 5, 56, 202, 1168, 190, [210, 226, 250, 236, 246])
headers = [
    ("Cluster", ""),
    ("Indicador 1", "meta / referência"),
    ("Indicador 2", "meta / referência"),
    ("Indicador 3", "meta / referência"),
    ("Indicador 4", "meta / referência"),
]
for col, (head, sub) in enumerate(headers):
    lines = [(head, {"size": 13, "bold": True, "color": WHITE})]
    if sub:
        lines.append((sub, {"size": 9.5, "color": "BFC9E3"}))
    set_cell(table.cell(0, col), lines, fill=NAVY)

# Linha de cluster "na meta" (todos verdes = OK)
row_a = [
    [("C.00", {"size": 16, "bold": True, "color": NAVY}), ("Cluster A", {"size": 10, "color": MUTED})],
    [("00,0%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("acima da meta", {"size": 10, "color": MUTED})],
    [("00,0%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("acima do esperado", {"size": 10, "color": MUTED})],
    [("00,0%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("acima da meta", {"size": 10, "color": MUTED})],
    [("0,00", {"size": 16, "bold": True, "color": GREEN_DARK}), ("dentro do limite", {"size": 10, "color": MUTED})],
]
# Linha de cluster "em risco" — mostra os 3 estados de cor: verde (OK), âmbar (ATENÇÃO), vermelho (RISCO)
row_b = [
    [("C.00", {"size": 16, "bold": True, "color": NAVY}), ("Cluster B", {"size": 10, "color": MUTED})],
    [("00,0%", {"size": 16, "bold": True, "color": GREEN_DARK}), ("na meta", {"size": 10, "color": MUTED})],
    [("00,0%", {"size": 16, "bold": True, "color": WARN}), ("no limite da meta", {"size": 10, "color": MUTED})],
    [("00,0%", {"size": 16, "bold": True, "color": RED}), ("abaixo da meta", {"size": 10, "color": MUTED})],
    [("0,00", {"size": 16, "bold": True, "color": RED}), ("acima do limite", {"size": 10, "color": MUTED})],
]
for col, lines in enumerate(row_a):
    set_cell(table.cell(1, col), lines, fill=WHITE)
for col, lines in enumerate(row_b):
    set_cell(table.cell(2, col), lines, fill="EAF2FB")

add_text(
    slide,
    "Como ler: [opcional] explique aqui, em uma frase, como interpretar o indicador menos óbvio "
    "da tabela. Se não precisar, apague esta caixa.",
    56, 402, 1130, 30, size=10.5, color=MUTED,
)

add_rect(slide, 56, 470, 545, 168, NAVY, round_corners=True)
add_text(slide, "DECISÃO OPERACIONAL", 80, 484, 320, 18, size=11, color="9FC9EE", bold=True)
add_text(
    slide,
    "[A decisão central da semana em 2 a 3 linhas: onde concentrar o esforço e por quê. "
    "É a frase que você quer que a diretoria leve como conclusão do slide.]",
    80, 506, 498, 118, size=14.5, color=WHITE, bold=True,
)

kpis = [
    ("[nº]", "[o que este número representa — evidência que sustenta a decisão]"),
    ("[nº]", "[segunda evidência de apoio]"),
    ("[nº]", "[terceira evidência — apague este cartão se usar só dois]"),
]
for index, (value, label) in enumerate(kpis):
    left = 625 + index * 202
    add_rect(slide, left, 470, 190, 168, WHITE, line=LINE, round_corners=True)
    add_text(slide, value, left + 15, 486, 162, 32, size=23, color=BLUE, bold=True)
    add_text(slide, label, left + 15, 522, 162, 104, size=11, color=MUTED)

add_notes(slide, (
    "SLIDE 1 — VISÃO EXECUTIVA\n"
    "\n"
    "1. Cabeçalho: troque REGIÃO, CÓDIGO e DATA na tarja azul; escreva a manchete (título grande) "
    "e o resumo (texto cinza) logo abaixo.\n"
    "2. Tabela de indicadores: renomeie os 4 indicadores no cabeçalho azul-marinho e preencha os "
    "valores dos clusters. A COR DO NÚMERO é o semáforo:\n"
    "   • VERDE = na meta   • ÂMBAR/DOURADO = atenção, no limite   • VERMELHO = abaixo da meta.\n"
    "   Para trocar a cor: selecione o número > Página Inicial > Cor da Fonte. A 1ª linha está toda "
    "verde e a 2ª mostra os três estados, de propósito, como referência.\n"
    "   Para adicionar/remover um cluster: clique na tabela > guia Layout > Inserir/Excluir linha.\n"
    "3. 'Como ler': caixa opcional. Apague se o indicador for autoexplicativo.\n"
    "4. Bloco azul-marinho = a decisão. Os 3 cartões à direita = evidências que sustentam a decisão "
    "(pode usar 2 ou 3; para 2, apague o último e centralize os que sobrarem)."
))

# ================================================================ Slide 2
slide = prs.slides.add_slide(blank)
chrome(slide, 2, 4, "Fonte: [Analítico Nominal e painel de improdutividade do período]")
title_block(
    slide,
    "DIAGNÓSTICO DE PRODUÇÃO",
    "[Diagnóstico em uma frase: onde está concentrado o problema de produção]",
    title_size=28,
)

# painel esquerdo — distribuição de produtividade
add_rect(slide, 56, 150, 610, 372, WHITE, line=LINE, round_corners=True)
add_text(slide, "[00] técnicos elegíveis", 78, 166, 330, 20, size=15, color=NAVY, bold=True)
add_text(slide, "[00 abaixo da meta (00%)]", 380, 168, 266, 18, size=13, color=RED, bold=True, align=PP_ALIGN.RIGHT)
add_text(
    slide,
    "[Critério da distribuição: por quartil de produtividade ou por tempo de casa. Explique a régua em uma linha.]",
    78, 190, 566, 16, size=10, color=MUTED,
)
# Barra empilhada — as larguras somam 100%. Ajuste cada segmento em Formatar Forma > Tamanho > Largura.
segments = [(40, GREEN_DARK), (25, AMBER), (25, ORANGE), (10, RED)]
bar_left, bar_width = 78, 566
offset = bar_left
for pct, color in segments:
    seg_width = bar_width * pct / 100
    add_rect(slide, offset, 214, seg_width, 20, color)
    offset += seg_width
legend = [
    ("40%", "[faixa mais produtiva]", GREEN_DARK),
    ("25%", "[faixa intermediária]", AMBER),
    ("25%", "[abaixo da intermediária]", ORANGE),
    ("10%", "[mais baixa / sem produção]", RED),
]
for index, (count, label, color) in enumerate(legend):
    left = 78 + index * 144
    add_rect(slide, left, 244, 7, 28, color)
    add_text(slide, [(count, {"size": 14, "bold": True}), (label, {"size": 8.5, "color": MUTED})],
             left + 12, 242, 132, 34)
add_rect(slide, 78, 284, 566, 42, "FDF4EA")
add_rect(slide, 78, 284, 5, 42, ORANGE)
add_text(
    slide,
    "[Leitura do quadro: a conclusão que a distribuição acima permite. Ex.: o risco não é o time "
    "inteiro, está concentrado em poucos nomes.]",
    92, 290, 544, 32, size=10.5, color="5A4632",
)
add_text(slide, "Gap estimado por cluster", 78, 340, 380, 18, size=13, color=NAVY, bold=True)
add_text(slide, "[00 OS no período]", 430, 341, 216, 16, size=12, color=RED, bold=True, align=PP_ALIGN.RIGHT)
hbar(slide, "C.00 Cluster B", "00", 78, 368, 130, 380, 1.0)
hbar(slide, "C.00 Cluster A", "00", 78, 392, 130, 380, 0.55, color="78A9D8")
add_text(
    slide,
    "[Como o gap é calculado: defina a fórmula em uma linha para dar rastreabilidade ao número.]",
    78, 420, 566, 30, size=9.5, color=MUTED,
)

# painel direito — motivos de improdutividade
add_rect(slide, 686, 150, 538, 372, WHITE, line=LINE, round_corners=True)
add_text(slide, "[000] visitas improdutivas", 708, 166, 320, 20, size=15, color=NAVY, bold=True)
add_text(slide, "[000 com motivo registrado]", 926, 168, 278, 18, size=12, color=RED, bold=True, align=PP_ALIGN.RIGHT)
add_text(
    slide,
    "Motivos registrados nos 4 maiores grupos ([período]):",
    708, 190, 494, 30, size=10, color=MUTED,
)
# Larguras das barras proporcionais ao maior motivo (o maior = 100%).
causes = [
    ("[Motivo 1 — o mais frequente]", "00", 1.0),
    ("[Motivo 2]", "00", 0.80),
    ("[Motivo 3]", "00", 0.60),
    ("[Motivo 4]", "00", 0.45),
]
for index, (label, value, ratio) in enumerate(causes):
    hbar(slide, label, value, 708, 236 + index * 34, 210, 220, ratio, color=ORANGE)
add_rect(slide, 708, 380, 494, 92, "F1F9F1")
add_rect(slide, 708, 380, 5, 92, GREEN)
add_text(slide, "00%", 724, 392, 90, 30, size=19, color=GREEN_DARK, bold=True)
add_text(
    slide,
    "[Insight sobre os motivos: ex.: parte relevante nasce antes do despacho, então cobrar só o "
    "técnico não resolve.]",
    822, 388, 368, 44, size=10.5, color="4A5B52",
)
add_text(
    slide,
    "[Opcional: observação adicional, como baixas sem motivo registrado a reconciliar com o BI.]",
    724, 434, 466, 32, size=9, color="4A5B52",
)

solutions = [
    ("01", "[Frente de ação 1]",
     "[Descreva a primeira frente de recuperação em 2 a 3 linhas: o que fazer, com que ritmo e até quando.]"),
    ("02", "[Frente de ação 2]",
     "[Segunda frente — normalmente o plano individual nominal, detalhado no slide 3.]"),
    ("03", "[Frente de ação 3]",
     "[Terceira frente — normalmente qualidade da demanda: travas de abertura e confirmação de cliente.]"),
]
for index, (number, head, body) in enumerate(solutions):
    left = 56 + index * 398
    add_rect(slide, left, 540, 382, 118, WASH)
    add_rect(slide, left, 540, 382, 3, BLUE)
    add_text(slide, number, left + 14, 552, 40, 26, size=20, color="A2B2C6", bold=True)
    add_text(slide, head, left + 54, 552, 316, 20, size=13, color=NAVY, bold=True)
    add_text(slide, body, left + 54, 574, 316, 78, size=10, color=MUTED)

add_notes(slide, (
    "SLIDE 2 — DIAGNÓSTICO DE PRODUÇÃO\n"
    "\n"
    "PAINEL ESQUERDO (distribuição da produtividade):\n"
    "• A barra colorida é empilhada: os 4 segmentos devem somar a largura total. Para mudar a "
    "proporção, clique num segmento e ajuste a LARGURA em Formatar Forma > Tamanho (ou arraste a "
    "alça lateral). Reposicione o segmento seguinte para encostar no anterior.\n"
    "• Atualize os percentuais e rótulos da legenda logo abaixo.\n"
    "• Barras de gap por cluster: a maior barra = 100%; as outras, proporcionais. Ajuste largura em "
    "Formatar Forma.\n"
    "\n"
    "PAINEL DIREITO (motivos de improdutividade):\n"
    "• Liste os 4 maiores motivos. A barra do 1º é a mais longa (100%); as demais, proporcionais.\n"
    "• A caixa verde traz o insight-chave sobre os motivos.\n"
    "\n"
    "RODAPÉ: as 3 frentes de ação (01/02/03). Mantenha títulos curtos e a descrição objetiva."
))

# ================================================================ Slide 3
slide = prs.slides.add_slide(blank)
chrome(slide, 3, 4, "Fonte: [Analítico Nominal por técnico, período] · [00] técnicos elegíveis · gap estimado")
title_block(
    slide,
    "GESTÃO NOMINAL · [DD A DD/MM/AAAA]",
    "[A recuperação tem nome: X áreas e Y técnicos concentram a maior parte do gap]",
    title_size=28,
)

add_rect(slide, 56, 148, 452, 400, WHITE, line=LINE, round_corners=True)
add_text(slide, "Gap por área de gestão (GA)", 78, 164, 300, 20, size=14, color=NAVY, bold=True)
add_text(slide, "em OS não executadas (estimado)", 300, 166, 188, 18, size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)
gas = [
    ("[Gestor 1] · C.00", "00", 1.0, BLUE),
    ("[Gestor 2] · C.00", "00", 0.85, BLUE),
    ("[Gestor 3] · C.00", "00", 0.70, BLUE),
    ("[Gestor 4] · C.00", "00", 0.55, "78A9D8"),
    ("[Gestor 5] · C.00", "00", 0.40, "78A9D8"),
]
for index, (label, value, ratio, color) in enumerate(gas):
    hbar(slide, label, value, 78, 202 + index * 32, 172, 180, ratio, color=color)
add_rect(slide, 78, 380, 408, 150, "F1F9F1")
add_rect(slide, 78, 380, 5, 150, GREEN)
add_text(slide, "[00%]", 94, 392, 120, 30, size=20, color=GREEN_DARK, bold=True)
add_text(
    slide,
    "[Destaque de apoio da coluna de GA: ex.: quanto da perda total está na área líder, ou um "
    "técnico sem nenhuma baixa no período e o que investigar.]",
    94, 424, 376, 96, size=10.5, color="4A5B52",
)

table = add_table(slide, 4, 3, 528, 148, 696, 400, [232, 240, 224])
set_cell(table.cell(0, 0), [("TÉCNICOS", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
set_cell(table.cell(0, 1), [("O QUE OS DADOS MOSTRAM", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
set_cell(table.cell(0, 2), [("AÇÃO IMEDIATA", {"size": 9.5, "bold": True, "color": WHITE})], fill=NAVY, v_margin=4)
nominal_rows = [
    (
        [("[Grupo de casos 1]", {"size": 10, "bold": True, "color": NAVY}),
         ("[nomes dos técnicos deste grupo]", {"size": 8, "color": MUTED})],
        [("[O que os dados mostram sobre este grupo, em 1 a 2 frases.]", {"size": 8.5})],
        [("[Ação imediata e prazo para este grupo.]", {"size": 8.5})],
        WHITE,
    ),
    (
        [("[Grupo de casos 2 — ex.: 8 maiores desvios]", {"size": 10, "bold": True, "color": NAVY}),
         ("[nomes e gaps individuais]", {"size": 8, "color": MUTED})],
        [("[Diagnóstico do grupo: quanto concentram do gap, contexto de escala etc.]", {"size": 8.5})],
        [("[Ação: conversa individual, plano de recuperação, verificação de escala.]", {"size": 8.5})],
        "F8FAFD",
    ),
    (
        [("Referências a replicar", {"size": 10, "bold": True, "color": GREEN_DARK}),
         ("[técnicos acima da meta a usar como padrão]", {"size": 8, "color": MUTED})],
        [("[Por que servem de referência: produtividade acima da meta nas mesmas condições.]", {"size": 8.5})],
        [("[Ação: mapear rota, carga e método e padronizar nos planos individuais.]", {"size": 8.5})],
        "F1F9F1",
    ),
]
for row_index, (col_a, col_b, col_c, fill) in enumerate(nominal_rows, start=1):
    set_cell(table.cell(row_index, 0), col_a, fill=fill, v_margin=4)
    set_cell(table.cell(row_index, 1), col_b, fill=fill, v_margin=4)
    set_cell(table.cell(row_index, 2), col_c, fill=fill, v_margin=4)
for row in table.rows:
    row.height = px(100)

add_notes(slide, (
    "SLIDE 3 — GESTÃO NOMINAL\n"
    "\n"
    "PAINEL ESQUERDO (gap por área de gestão / GA):\n"
    "• Liste os gestores com maior gap. A maior barra = 100%; as demais, proporcionais (ajuste a "
    "largura em Formatar Forma). As duas últimas usam azul mais claro só para dar hierarquia visual.\n"
    "• A caixa verde é um destaque de apoio (percentual concentrado, técnico parado etc.).\n"
    "\n"
    "TABELA À DIREITA (casos que exigem decisão individual):\n"
    "• Cada linha é um grupo de casos: coluna 1 = quem, coluna 2 = o que os dados mostram, coluna 3 = "
    "ação imediata.\n"
    "• A ÚLTIMA linha (verde) é para 'Referências a replicar' — técnicos acima da meta. Se não quiser "
    "usar, exclua a linha (clique na tabela > Layout > Excluir linha).\n"
    "• Para mais grupos, adicione linhas pela guia Layout."
))

# ================================================================ Slide 4
slide = prs.slides.add_slide(blank)
chrome(slide, 4, 4, "Causas não comprovadas pelos dados permanecem marcadas como “em validação”")
title_block(
    slide,
    "FCA · FATO, CAUSA E AÇÃO",
    "[Quantas frentes colocam a regional em recuperação nesta semana]",
    title_size=28,
)

table = add_table(slide, 5, 6, 42, 148, 1196, 452, [200, 258, 268, 116, 216, 138])
fca_heads = ["FATO", "CAUSA", "AÇÃO", "PRAZO", "RESPONSÁVEL", "STATUS"]
for col, head in enumerate(fca_heads):
    set_cell(table.cell(0, col), [(head, {"size": 10.5, "bold": True, "color": WHITE})], fill=NAVY)

# As 4 linhas mostram os 4 estilos de STATUS disponíveis — use o que couber em cada caso.
fca_rows = [
    (
        [("[Fato 1 — o problema observado]", {"size": 11, "bold": True, "color": NAVY}),
         ("[dado que quantifica o fato]", {"size": 9, "color": MUTED})],
        [("[Causa: por que aconteceu, com o que os dados sustentam.]", {"size": 9.5})],
        [("[Ação concreta para tratar a causa.]", {"size": 9.5})],
        [("[DD/MM]", {"size": 11, "bold": True, "color": NAVY}), ("[revisão]", {"size": 9, "color": MUTED})],
        [("[Responsável]", {"size": 9.5})],
        ("A iniciar", BLUE_LIGHT, "114F88"),
        WHITE,
    ),
    (
        [("[Fato 2]", {"size": 11, "bold": True, "color": NAVY}),
         ("[dado de apoio]", {"size": 9, "color": MUTED})],
        [("[Causa.]", {"size": 9.5})],
        [("[Ação.]", {"size": 9.5})],
        [("[DD/MM]", {"size": 11, "bold": True, "color": NAVY}), ("[planos]", {"size": 9, "color": MUTED})],
        [("[Responsável]", {"size": 9.5})],
        ("Imediato", RED_LIGHT, "8A1717"),
        "F8FAFD",
    ),
    (
        [("[Fato 3]", {"size": 11, "bold": True, "color": NAVY}),
         ("[dado de apoio]", {"size": 9, "color": MUTED})],
        [("[Causa ainda não comprovada pelos dados.]", {"size": 9.5})],
        [("[Ação, incluindo o que falta para comprovar a causa.]", {"size": 9.5})],
        [("[DD/MM]", {"size": 11, "bold": True, "color": NAVY}), ("[travas]", {"size": 9, "color": MUTED})],
        [("[Responsável]", {"size": 9.5})],
        ("Em validação", "FFF0C7", "765400"),
        WHITE,
    ),
    (
        [("[Fato 4]", {"size": 11, "bold": True, "color": NAVY}),
         ("[dado de apoio]", {"size": 9, "color": MUTED})],
        [("[Causa.]", {"size": 9.5})],
        [("[Ação.]", {"size": 9.5})],
        [("[DD/MM]", {"size": 11, "bold": True, "color": NAVY}), ("[meta]", {"size": 9, "color": MUTED})],
        [("[Responsável]", {"size": 9.5})],
        ("Atenção", "FFE5CB", "8A4C00"),
        "F8FAFD",
    ),
]
for row_index, (fato, causa, acao, prazo, resp, status, fill) in enumerate(fca_rows, start=1):
    set_cell(table.cell(row_index, 0), fato, fill=fill)
    set_cell(table.cell(row_index, 1), causa, fill=fill)
    set_cell(table.cell(row_index, 2), acao, fill=fill)
    set_cell(table.cell(row_index, 3), prazo, fill=fill)
    set_cell(table.cell(row_index, 4), resp, fill=fill)
    status_text, status_fill, status_color = status
    set_cell(
        table.cell(row_index, 5),
        [(status_text, {"size": 10, "bold": True, "color": status_color, "align": PP_ALIGN.CENTER})],
        fill=status_fill,
    )

add_rect(slide, 42, 618, 1196, 48, NAVY)
add_rect(slide, 42, 618, 7, 48, GREEN)
add_text(slide, "DECISÕES SOLICITADAS", 64, 634, 200, 18, size=11.5, color=WHITE, bold=True)
add_text(
    slide,
    "[Decisão 1]   ·   [Decisão 2]   ·   [Decisão 3]",
    270, 634, 940, 18, size=11.5, color="D5DEF3",
)

add_notes(slide, (
    "SLIDE 4 — FCA (FATO, CAUSA, AÇÃO)\n"
    "\n"
    "• Uma linha por frente: Fato (o problema + dado) / Causa / Ação / Prazo / Responsável / Status.\n"
    "• A coluna STATUS tem 4 estilos prontos, um por linha, como amostra — use o que couber:\n"
    "   A INICIAR (azul) · IMEDIATO (vermelho) · EM VALIDAÇÃO (amarelo) · ATENÇÃO (laranja).\n"
    "   Para mudar o estilo de uma célula: selecione-a, troque o preenchimento (Formatar Forma > "
    "Preenchimento) e a cor da fonte, ou copie a célula da linha que já tem o estilo desejado.\n"
    "• Regra editorial: causa não comprovada pelos dados fica como 'Em validação'.\n"
    "• Para menos frentes, exclua linhas (Layout > Excluir linha) e estique as demais.\n"
    "• Barra azul-marinho no rodapé = decisões que você pede à diretoria (separadas por · )."
))

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(f"OK: {OUTPUT}")
